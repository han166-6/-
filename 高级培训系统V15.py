"""
客户经理培训档案可视化系统 - 增强版（含多浏览器截图导出功能）
包含数据分析、报告生成、预警提醒、能力标签、经典案例功能
支持Chrome/Firefox/Edge截图导出PDF，完美保留所有图表和样式
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
from datetime import datetime, timedelta
import warnings
import io
from io import BytesIO
import base64
import random
import json
import math
import os
import tempfile
from PIL import Image
import time
import subprocess
import sys
import socket
import pickle  # 新增：用于数据持久化

# 尝试导入导出功能所需的库
try:
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# 尝试导入截图相关库
try:
    from selenium import webdriver
    from selenium.webdriver.common.by import By
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    from webdriver_manager.chrome import ChromeDriverManager
    from webdriver_manager.firefox import GeckoDriverManager
    from webdriver_manager.microsoft import EdgeChromiumDriverManager
    
    try:
        from selenium.webdriver.chrome.service import Service as ChromeService
        from selenium.webdriver.chrome.options import Options as ChromeOptions
        CHROME_SUPPORT = True
    except:
        CHROME_SUPPORT = False
    
    try:
        from selenium.webdriver.firefox.service import Service as FirefoxService
        from selenium.webdriver.firefox.options import Options as FirefoxOptions
        FIREFOX_SUPPORT = True
    except:
        FIREFOX_SUPPORT = False
    
    try:
        from selenium.webdriver.edge.service import Service as EdgeService
        from selenium.webdriver.edge.options import Options as EdgeOptions
        EDGE_SUPPORT = True
    except:
        EDGE_SUPPORT = False
    
    SELENIUM_AVAILABLE = True
except ImportError:
    SELENIUM_AVAILABLE = False
    CHROME_SUPPORT = FIREFOX_SUPPORT = EDGE_SUPPORT = False

warnings.filterwarnings('ignore')

# 设置页面配置
st.set_page_config(
    page_title="客户经理培训档案系统",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# 自定义CSS美化界面
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        color: #1E3A8A;
        font-weight: bold;
        margin-bottom: 1rem;
        text-align: center;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    .metric-card {
        background: white;
        border-radius: 12px;
        padding: 20px;
        text-align: center;
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.08);
        border: 1px solid #e5e7eb;
        transition: transform 0.3s;
    }
    .metric-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 8px 16px rgba(0, 0, 0, 0.12);
    }
    .section-title {
        font-size: 1.5rem;
        color: #1E3A8A;
        font-weight: bold;
        margin: 2rem 0 1rem 0;
        padding-bottom: 0.5rem;
        border-bottom: 3px solid #3B82F6;
    }
    .subsection-title {
        font-size: 1.2rem;
        color: #374151;
        font-weight: 600;
        margin: 1.5rem 0 0.8rem 0;
        padding-left: 10px;
        border-left: 4px solid #10B981;
    }
    .warning-box {
        background: linear-gradient(135deg, #fef3c7 0%, #fde68a 100%);
        border-left: 4px solid #F59E0B;
        padding: 15px;
        border-radius: 8px;
        margin: 10px 0;
    }
    .warning-box.high {
        background: linear-gradient(135deg, #fecaca 0%, #fca5a5 100%);
        border-left: 4px solid #EF4444;
    }
    .warning-box.medium {
        background: linear-gradient(135deg, #fef3c7 0%, #fde68a 100%);
        border-left: 4px solid #F59E0B;
    }
    .warning-box.low {
        background: linear-gradient(135deg, #dbeafe 0%, #bfdbfe 100%);
        border-left: 4px solid #3B82F6;
    }
    .success-box {
        background: linear-gradient(135deg, #d1fae5 0%, #a7f3d0 100%);
        border-left: 4px solid #10B981;
        padding: 15px;
        border-radius: 8px;
        margin: 10px 0;
    }
    .tag-container {
        display: flex;
        flex-wrap: wrap;
        gap: 8px;
        margin: 10px 0;
    }
    .skill-tag {
        display: inline-flex;
        align-items: center;
        background: linear-gradient(135deg, #e0e7ff 0%, #c7d2fe 100%);
        color: #4F46E5;
        padding: 6px 14px;
        border-radius: 20px;
        font-size: 0.85rem;
        font-weight: 500;
        border: 1px solid #c7d2fe;
        transition: all 0.3s ease;
    }
    .skill-tag:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(79, 70, 229, 0.2);
    }
    .skill-tag.language {
        background: linear-gradient(135deg, #dbeafe 0%, #bfdbfe 100%);
        color: #1D4ED8;
        border-color: #93c5fd;
    }
    .skill-tag.management {
        background: linear-gradient(135deg, #fce7f3 0%, #fbcfe8 100%);
        color: #BE185D;
        border-color: #f9a8d4;
    }
    .skill-tag.technical {
        background: linear-gradient(135deg, #dcfce7 0%, #bbf7d0 100%);
        color: #166534;
        border-color: #86efac;
    }
    .skill-tag.certificate {
        background: linear-gradient(135deg, #fef3c7 0%, #fde68a 100%);
        color: #92400E;
        border-color: #fcd34d;
    }
    .weakness-tag {
        display: inline-flex;
        align-items: center;
        background: linear-gradient(135deg, #fee2e2 0%, #fecaca 100%);
        color: #DC2626;
        padding: 6px 14px;
        border-radius: 20px;
        font-size: 0.85rem;
        font-weight: 500;
        border: 1px solid #fca5a5;
        margin: 5px;
        transition: all 0.3s ease;
    }
    .weakness-tag:hover {
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(220, 38, 38, 0.2);
    }
    .team-weakness-card {
        background: linear-gradient(135deg, #fef2f2 0%, #fee2e2 100%);
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
        border-left: 4px solid #EF4444;
    }
    .team-weakness-title {
        font-size: 1.1rem;
        font-weight: 600;
        color: #991B1B;
        margin-bottom: 10px;
    }
    .team-weakness-item {
        display: inline-block;
        background: white;
        border-radius: 15px;
        padding: 5px 12px;
        margin: 3px;
        font-size: 0.9rem;
        border: 1px solid #FCA5A5;
    }
    .case-card {
        background: linear-gradient(135deg, #f8fafc 0%, #f1f5f9 100%);
        border-radius: 12px;
        padding: 20px;
        margin: 15px 0;
        border-left: 5px solid #3B82F6;
        transition: all 0.3s ease;
        position: relative;
        overflow: hidden;
    }
    .case-card::before {
        content: '';
        position: absolute;
        top: 0;
        left: 0;
        right: 0;
        height: 4px;
        background: linear-gradient(90deg, #3B82F6, #8B5CF6);
    }
    .case-card:hover {
        transform: translateY(-3px);
        box-shadow: 0 8px 25px rgba(0,0,0,0.1);
    }
    .case-title {
        font-size: 1.2rem;
        font-weight: 600;
        color: #1E3A8A;
        margin-bottom: 10px;
        display: flex;
        align-items: center;
        gap: 8px;
    }
    .case-content {
        color: #4B5563;
        line-height: 1.6;
        margin: 10px 0;
    }
    .case-meta {
        display: flex;
        gap: 15px;
        margin-top: 15px;
        padding-top: 15px;
        border-top: 1px solid #e5e7eb;
        color: #6B7280;
        font-size: 0.9rem;
    }
    .case-tag {
        background: linear-gradient(135deg, #e0e7ff 0%, #c7d2fe 100%);
        color: #4F46E5;
        padding: 4px 12px;
        border-radius: 15px;
        font-size: 0.8rem;
        font-weight: 500;
    }
    .person-card {
        background: white;
        border-radius: 12px;
        padding: 25px;
        margin: 20px 0;
        box-shadow: 0 4px 16px rgba(0,0,0,0.08);
        border: 1px solid #e5e7eb;
        transition: all 0.3s ease;
        position: relative;
    }
    .person-card:hover {
        transform: translateY(-5px);
        box-shadow: 0 12px 30px rgba(0,0,0,0.12);
    }
    .person-card::after {
        content: '';
        position: absolute;
        top: 0;
        left: 0;
        right: 0;
        height: 4px;
        background: linear-gradient(90deg, #3B82F6, #8B5CF6);
        border-radius: 12px 12px 0 0;
    }
    .stProgress > div > div > div > div {
        background: linear-gradient(90deg, #4F46E5 0%, #3B82F6 100%);
    }
    .team-stat-card {
        background: linear-gradient(135deg, #f8fafc 0%, #e5e7eb 100%);
        border-radius: 12px;
        padding: 20px;
        margin: 15px 0;
        border: 1px solid #d1d5db;
    }
    .photo-container {
        width: 120px;
        height: 160px;
        border-radius: 10px;
        overflow: hidden;
        box-shadow: 0 4px 12px rgba(0,0,0,0.15);
        margin: 0 auto;
        border: 3px solid white;
    }
    .photo-container img {
        width: 100%;
        height: 100%;
        object-fit: cover;
    }
    .icon {
        font-size: 1.2em;
        margin-right: 5px;
    }
    .dataframe {
        border-collapse: collapse;
        width: 100%;
    }
    .dataframe th {
        background-color: #3B82F6;
        color: white;
        padding: 10px;
        text-align: left;
    }
    .dataframe td {
        padding: 10px;
        border-bottom: 1px solid #ddd;
    }
    .dataframe tr:hover {
        background-color: #f5f5f5;
    }
    .comparison-container {
        background: linear-gradient(135deg, #f8fafc 0%, #e5e7eb 100%);
        border-radius: 15px;
        padding: 25px;
        margin: 20px 0;
    }
    .threshold-card {
        background: linear-gradient(135deg, #f8fafc 0%, #e5e7eb 100%);
        border-radius: 12px;
        padding: 20px;
        margin: 15px 0;
        border: 1px solid #3B82F6;
    }
    .export-section {
        background: linear-gradient(135deg, #f0f9ff 0%, #e6f0fa 100%);
        border-radius: 15px;
        padding: 25px;
        margin: 30px 0;
        border: 1px solid #3B82F6;
    }
    .diagnosis-box {
        background: #f8fafc;
        border: 1px solid #94a3b8;
        border-radius: 8px;
        padding: 15px;
        margin: 10px 0;
        font-family: monospace;
    }
</style>
""", unsafe_allow_html=True)

# 初始化Session State
if 'personnel_photos' not in st.session_state:
    st.session_state.personnel_photos = {}

if 'warning_thresholds' not in st.session_state:
    st.session_state.warning_thresholds = {
        '高级客户经理': 80,
        '客户经理': 70,
        '客户经理助理': 60,
        '其他': 70
    }

if 'team_threshold' not in st.session_state:
    st.session_state.team_threshold = 70

# 新增：用于数据持久化的状态
if 'data_loaded' not in st.session_state:
    st.session_state.data_loaded = False
if 'uploaded_file_name' not in st.session_state:
    st.session_state.uploaded_file_name = None
if 'analysis_data' not in st.session_state:
    st.session_state.analysis_data = None

# ========== 本地缓存持久化 ==========
CACHE_FILE = "training_analysis_cache.pkl"

def load_cached_data():
    """从本地缓存文件加载数据，返回 (analysis_data, file_name) 或 (None, None)"""
    if os.path.exists(CACHE_FILE):
        try:
            with open(CACHE_FILE, 'rb') as f:
                cached = pickle.load(f)
                return cached.get('analysis_data'), cached.get('file_name')
        except Exception as e:
            st.warning(f"读取缓存失败：{e}，将重新上传。")
            return None, None
    return None, None

def save_cached_data(analysis_data, file_name):
    """将数据保存到本地缓存文件"""
    try:
        with open(CACHE_FILE, 'wb') as f:
            pickle.dump({'analysis_data': analysis_data, 'file_name': file_name}, f)
    except Exception as e:
        st.error(f"保存缓存失败：{e}")

# ========== 截图导出功能 ==========
def check_browser_installed(browser_type='chrome'):
    """检查指定浏览器是否安装"""
    browser_type = browser_type.lower()
    
    common_paths = {
        'chrome': [
            r"C:\Program Files\Google\Chrome\Application\chrome.exe",
            r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
            os.path.expanduser(r"~\AppData\Local\Google\Chrome\Application\chrome.exe"),
        ],
        'firefox': [
            r"C:\Program Files\Mozilla Firefox\firefox.exe",
            r"C:\Program Files (x86)\Mozilla Firefox\firefox.exe",
            os.path.expanduser(r"~\AppData\Local\Mozilla Firefox\firefox.exe"),
        ],
        'edge': [
            r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
            r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
            os.path.expanduser(r"~\AppData\Local\Microsoft\Edge\Application\msedge.exe"),
        ]
    }
    
    paths = common_paths.get(browser_type, [])
    for path in paths:
        expanded_path = os.path.expanduser(path)
        if os.path.exists(expanded_path):
            return True
    
    try:
        if browser_type == 'chrome':
            if sys.platform == "win32":
                result = subprocess.run(['where', 'chrome'], capture_output=True, text=True)
            else:
                result = subprocess.run(['which', 'google-chrome'], capture_output=True, text=True)
        elif browser_type == 'firefox':
            if sys.platform == "win32":
                result = subprocess.run(['where', 'firefox'], capture_output=True, text=True)
            else:
                result = subprocess.run(['which', 'firefox'], capture_output=True, text=True)
        elif browser_type == 'edge':
            if sys.platform == "win32":
                result = subprocess.run(['where', 'msedge'], capture_output=True, text=True)
            else:
                result = subprocess.run(['which', 'microsoft-edge'], capture_output=True, text=True)
        else:
            return False
        return result.returncode == 0
    except:
        return False

def get_available_browsers():
    """获取系统上可用的浏览器列表"""
    available = []
    
    if check_browser_installed('chrome'):
        available.append(('Chrome', 'chrome'))
    if check_browser_installed('firefox'):
        available.append(('Firefox', 'firefox'))
    if check_browser_installed('edge'):
        available.append(('Edge', 'edge'))
    
    return available

def diagnose_connection():
    """诊断网络连接"""
    results = []
    
    test_hosts = [
        ("localhost", 8501),
        ("127.0.0.1", 8501),
        ("0.0.0.0", 8501)
    ]
    
    for host, port in test_hosts:
        try:
            sock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
            sock.settimeout(2)
            result = sock.connect_ex((host, port))
            sock.close()
            
            if result == 0:
                results.append(f"✅ {host}:{port} - 可以连接")
            else:
                results.append(f"❌ {host}:{port} - 无法连接 (错误码: {result})")
        except Exception as e:
            results.append(f"❌ {host}:{port} - 诊断失败: {str(e)}")
    
    return results

def capture_streamlit_screenshots(url, output_dir, browser_type='chrome', progress_callback=None):
    """使用 Selenium 截取 Streamlit 应用的各个页面（支持多浏览器）"""
    
    if not SELENIUM_AVAILABLE:
        raise ImportError("Selenium库未安装")
    
    browser_type = browser_type.lower()
    screenshot_paths = []
    driver = None
    max_retries = 3
    
    try:
        if browser_type == 'chrome':
            if not CHROME_SUPPORT:
                raise ImportError("Chrome支持库未完全安装")
            
            chrome_options = ChromeOptions()
            chrome_options.add_argument('--headless=new')
            chrome_options.add_argument('--no-sandbox')
            chrome_options.add_argument('--disable-dev-shm-usage')
            chrome_options.add_argument('--disable-gpu')
            chrome_options.add_argument('--window-size=1920,1080')
            chrome_options.add_argument('--disable-blink-features=AutomationControlled')
            chrome_options.add_argument('--ignore-certificate-errors')
            chrome_options.add_argument('--allow-insecure-localhost')
            
            if progress_callback:
                progress_callback("正在启动Chrome浏览器...")
            
            service = ChromeService(ChromeDriverManager().install())
            driver = webdriver.Chrome(service=service, options=chrome_options)
            
        elif browser_type == 'firefox':
            if not FIREFOX_SUPPORT:
                raise ImportError("Firefox支持库未完全安装")
            
            firefox_options = FirefoxOptions()
            firefox_options.add_argument('--headless')
            firefox_options.add_argument('--width=1920')
            firefox_options.add_argument('--height=1080')
            firefox_options.add_argument('--ignore-certificate-errors')
            
            if progress_callback:
                progress_callback("正在启动Firefox浏览器...")
            
            service = FirefoxService(GeckoDriverManager().install())
            driver = webdriver.Firefox(service=service, options=firefox_options)
            
        elif browser_type == 'edge':
            if not EDGE_SUPPORT:
                raise ImportError("Edge支持库未完全安装")
            
            edge_options = EdgeOptions()
            edge_options.add_argument('--headless=new')
            edge_options.add_argument('--no-sandbox')
            edge_options.add_argument('--disable-gpu')
            edge_options.add_argument('--window-size=1920,1080')
            edge_options.add_argument('--ignore-certificate-errors')
            edge_options.add_argument('--allow-insecure-localhost')
            
            if progress_callback:
                progress_callback("正在启动Edge浏览器...")
            
            service = EdgeService(EdgeChromiumDriverManager().install())
            driver = webdriver.Edge(service=service, options=edge_options)
        else:
            raise ValueError(f"不支持的浏览器类型: {browser_type}")
        
        # 设置页面加载超时
        driver.set_page_load_timeout(30)
        
        if progress_callback:
            progress_callback(f"正在打开应用: {url}...")
        
        # 尝试加载页面，带重试
        for attempt in range(max_retries):
            try:
                driver.get(url)
                break
            except Exception as e:
                if attempt == max_retries - 1:
                    raise
                if progress_callback:
                    progress_callback(f"连接失败，重试 {attempt + 1}/{max_retries}...")
                time.sleep(2)
        
        time.sleep(3)
        
        # 检查页面是否成功加载
        try:
            WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.CSS_SELECTOR, '[data-baseweb="tab"]'))
            )
        except:
            # 如果没有找到tab，可能是页面加载失败
            page_source = driver.page_source
            if "streamlit" not in page_source.lower():
                raise Exception("页面加载失败：不是有效的Streamlit应用")
        
        # 获取所有 tab
        tabs = driver.find_elements(By.CSS_SELECTOR, '[data-baseweb="tab"]')
        
        # 如果没有找到tab，尝试其他选择器
        if not tabs:
            tabs = driver.find_elements(By.CSS_SELECTOR, '.stTabs [role="tab"]')
        
        if not tabs:
            # 如果还是没有tab，说明页面结构不同，只截取当前页面
            if progress_callback:
                progress_callback("未找到标签页，将截取当前页面...")
            
            # 滚动到顶部
            driver.execute_script("window.scrollTo(0, 0);")
            time.sleep(0.5)
            
            # 获取页面高度并截图
            total_height = driver.execute_script("return document.body.scrollHeight")
            driver.set_window_size(1920, total_height)
            time.sleep(0.5)
            
            # 截图
            screenshot_path = os.path.join(output_dir, "full_page.png")
            driver.save_screenshot(screenshot_path)
            screenshot_paths.append(screenshot_path)
        else:
            tab_names = ['数据概览', '个人档案', '对比分析', '团队分析', '数据报告']
            
            for i, tab_name in enumerate(tab_names):
                if i < len(tabs):
                    if progress_callback:
                        progress_callback(f"正在截取: {tab_name}...")
                    
                    try:
                        driver.execute_script("arguments[0].click();", tabs[i])
                        time.sleep(2)
                        
                        WebDriverWait(driver, 10).until(
                            EC.presence_of_element_located((By.CLASS_NAME, "stApp"))
                        )
                        time.sleep(1)
                        
                        driver.execute_script("window.scrollTo(0, 0);")
                        time.sleep(0.5)
                        
                        total_height = driver.execute_script("return document.body.scrollHeight")
                        driver.set_window_size(1920, total_height)
                        time.sleep(0.5)
                        
                        screenshot_path = os.path.join(output_dir, f"tab_{i+1}_{tab_name}.png")
                        driver.save_screenshot(screenshot_path)
                        screenshot_paths.append(screenshot_path)
                    except Exception as e:
                        if progress_callback:
                            progress_callback(f"截取 {tab_name} 时出错: {str(e)}")
                        continue
        
        if progress_callback:
            progress_callback(f"截图完成！共 {len(screenshot_paths)} 张")
        
        return screenshot_paths
        
    except Exception as e:
        raise e
    finally:
        if driver:
            driver.quit()

def combine_screenshots_to_pdf(screenshot_paths, output_pdf_path, progress_callback=None):
    """将多个截图合并为 PDF"""
    
    if not REPORTLAB_AVAILABLE:
        raise ImportError("ReportLab库未安装")
    
    c = canvas.Canvas(output_pdf_path, pagesize=A4)
    
    for i, screenshot_path in enumerate(screenshot_paths):
        if progress_callback:
            progress_callback(f"正在生成PDF: 第 {i+1}/{len(screenshot_paths)} 页...")
        
        if i > 0:
            c.showPage()
        
        img = ImageReader(screenshot_path)
        img_width, img_height = img.getSize()
        
        page_width, page_height = A4
        margin = 36
        max_width = page_width - 2 * margin
        max_height = page_height - 2 * margin
        
        scale = min(max_width / img_width, max_height / img_height)
        
        x = (page_width - img_width * scale) / 2
        y = (page_height - img_height * scale) / 2
        
        c.drawImage(img, x, y, width=img_width * scale, height=img_height * scale)
    
    c.save()
    
    if progress_callback:
        progress_callback("PDF生成完成！")

def create_pdf_from_screenshots(url, browser_type='chrome', progress_callback=None):
    """主函数：从 Streamlit 截图创建 PDF"""
    
    # 尝试多个可能的本地地址
    possible_urls = [
        "http://localhost:8501",
        "http://127.0.0.1:8501",
        "http://0.0.0.0:8501"
    ]
    
    # 如果传入的url在possible_urls中，就使用它
    if url in possible_urls:
        possible_urls = [url] + [u for u in possible_urls if u != url]
    
    last_error = None
    
    for test_url in possible_urls:
        try:
            if progress_callback:
                progress_callback(f"尝试连接: {test_url}...")
            
            with tempfile.TemporaryDirectory() as tmpdir:
                screenshot_paths = capture_streamlit_screenshots(
                    test_url, tmpdir, browser_type, progress_callback
                )
                
                if screenshot_paths:
                    pdf_path = os.path.join(tmpdir, "output.pdf")
                    combine_screenshots_to_pdf(screenshot_paths, pdf_path, progress_callback)
                    
                    with open(pdf_path, 'rb') as f:
                        pdf_bytes = f.read()
                    
                    return pdf_bytes
                    
        except Exception as e:
            last_error = e
            if progress_callback:
                progress_callback(f"连接 {test_url} 失败，尝试下一个...")
            continue
    
    if last_error:
        raise Exception(f"无法连接到Streamlit应用。请确保：\n1. 应用正在运行\n2. 防火墙没有阻止连接\n3. 尝试使用不同的浏览器\n\n详细错误：{str(last_error)}")
    
    return None

def install_selenium():
    """安装Selenium相关库"""
    try:
        st.info("正在安装 Selenium 和 webdriver-manager...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "selenium", "webdriver-manager"])
        st.success("✅ 安装成功！请重新运行应用。")
        st.info("请点击浏览器上的刷新按钮重新加载页面。")
    except Exception as e:
        st.error(f"安装失败: {str(e)}")
        st.info("请手动运行以下命令安装：")
        st.code("pip install selenium webdriver-manager")

# ========== 原始功能函数 ==========
def load_excel_file(uploaded_file):
    """加载Excel文件"""
    try:
        xls = pd.ExcelFile(uploaded_file)
        
        required_sheets = ['人员信息', '培训课程', '能力评估', '岗位培训要求']
        missing_sheets = [sheet for sheet in required_sheets if sheet not in xls.sheet_names]
        
        if missing_sheets:
            st.error(f"❌ 缺少必需的工作表: {', '.join(missing_sheets)}")
            return None
        
        personnel_df = pd.read_excel(xls, sheet_name='人员信息')
        training_df = pd.read_excel(xls, sheet_name='培训课程')
        capability_df = pd.read_excel(xls, sheet_name='能力评估')
        training_requirements_df = pd.read_excel(xls, sheet_name='岗位培训要求')
        
        return {
            '人员信息': personnel_df,
            '培训课程': training_df,
            '能力评估': capability_df,
            '岗位培训要求': training_requirements_df
        }
    except Exception as e:
        st.error(f"❌ 读取Excel文件失败: {str(e)}")
        return None

def parse_skill_tags(tags_str):
    """解析能力标签字符串"""
    if pd.isna(tags_str) or not tags_str or str(tags_str).strip() == '':
        return []
    
    tags_str = str(tags_str).strip()
    tags = []
    
    tags_str = tags_str.replace('，', ',')
    tag_list = [tag.strip() for tag in tags_str.split(',') if tag.strip()]
    
    if not tag_list:
        for separator in [';', '、', '|', '/', ' ']:
            if separator in tags_str:
                tag_list = [tag.strip() for tag in tags_str.split(separator) if tag.strip()]
                break
    
    if not tag_list:
        tag_list = [tags_str]
    
    return list(dict.fromkeys(tag_list))

def parse_classic_cases(case_str):
    """解析经典案例字符串"""
    if pd.isna(case_str) or not case_str or str(case_str).strip() == '':
        return []
    
    case_str = str(case_str).strip()
    cases = []
    
    if case_str.startswith('[') and case_str.endswith(']'):
        try:
            cases_data = json.loads(case_str)
            for case in cases_data:
                if isinstance(case, dict):
                    cases.append({
                        'title': case.get('title', '经典案例'),
                        'content': case.get('content', ''),
                        'year': str(case.get('year', '')),
                        'type': case.get('type', '项目案例'),
                        'role': case.get('role', ''),
                        'achievement': case.get('achievement', '')
                    })
            return cases
        except:
            pass
    
    lines = [line.strip() for line in case_str.split('\n') if line.strip()]
    
    for line in lines:
        case_data = {
            'title': '经典案例',
            'content': '',
            'year': '',
            'type': '项目案例',
            'role': '',
            'achievement': ''
        }
        
        if '|' in line:
            parts = [p.strip() for p in line.split('|')]
            if len(parts) >= 1:
                case_data['title'] = parts[0]
            if len(parts) >= 2:
                case_data['content'] = parts[1]
            if len(parts) >= 3:
                case_data['year'] = parts[2]
            if len(parts) >= 4:
                case_data['type'] = parts[3]
            if len(parts) >= 5:
                case_data['role'] = parts[4]
            if len(parts) >= 6:
                case_data['achievement'] = parts[5]
        elif '，' in line:
            parts = [p.strip() for p in line.split('，')]
            if len(parts) >= 1:
                case_data['title'] = parts[0]
            if len(parts) >= 2:
                case_data['content'] = parts[1]
            if len(parts) >= 3:
                case_data['year'] = parts[2]
            if len(parts) >= 4:
                case_data['type'] = parts[3]
            if len(parts) >= 5:
                case_data['role'] = parts[4]
            if len(parts) >= 6:
                case_data['achievement'] = parts[5]
        elif ',' in line:
            parts = [p.strip() for p in line.split(',')]
            if len(parts) >= 1:
                case_data['title'] = parts[0]
            if len(parts) >= 2:
                case_data['content'] = parts[1]
            if len(parts) >= 3:
                case_data['year'] = parts[2]
            if len(parts) >= 4:
                case_data['type'] = parts[3]
            if len(parts) >= 5:
                case_data['role'] = parts[4]
            if len(parts) >= 6:
                case_data['achievement'] = parts[5]
        else:
            case_data['title'] = line
            case_data['content'] = line
        
        cases.append(case_data)
    
    return cases

def get_tag_category(tag):
    """获取标签分类"""
    tag_lower = tag.lower()
    
    language_keywords = ['英语', '日语', '法语', '德语', '俄语', '西班牙语', '韩语', '阿拉伯语', '语言', '口语', '翻译', '专八', '六级', '四级']
    if any(keyword in tag_lower for keyword in language_keywords):
        return 'language'
    
    management_keywords = ['管理', '领导', '团队', '沟通', '协调', '组织', '规划', '战略', '决策', '谈判', '团队管理', '项目管理']
    if any(keyword in tag_lower for keyword in management_keywords):
        return 'management'
    
    certificate_keywords = ['证书', '认证', '资格', '执照', 'pmp', 'cfa', 'cpa', '认证证书']
    if any(keyword in tag_lower for keyword in certificate_keywords):
        return 'certificate'
    
    technical_keywords = ['技术', '分析', '数据', '编程', '开发', '设计', '工程', '研发', '测试', '运维', '数据分析', '技术开发']
    if any(keyword in tag_lower for keyword in technical_keywords):
        return 'technical'
    
    return 'default'

def get_case_icon(case_type):
    """根据案例类型获取图标"""
    icons = {
        '项目案例': '📋',
        '技术支持': '🛠️',
        '客户服务': '💼',
        '问题解决': '✅',
        '创新成果': '💡',
        '团队建设': '👥',
        '培训指导': '🎓',
        '其他': '📝'
    }
    return icons.get(case_type, '📋')

def get_warning_level(score, threshold):
    """根据分数和阈值获取预警级别"""
    if score < threshold * 0.7:
        return 'high'
    elif score < threshold * 0.85:
        return 'medium'
    elif score < threshold:
        return 'low'
    else:
        return None

def calculate_training_completion_by_position(person, position, category_stats, position_requirements):
    """根据岗位要求计算培训完成度"""
    requirements = position_requirements.get(position, {})
    
    if not requirements:
        total_required = 10
        total_completed = 0
        for category, stats in category_stats.items():
            total_completed += stats.get('completed', 0) + stats.get('assessed', 0)
        overall_rate = (total_completed / total_required * 100) if total_required > 0 else 0
        
        return {
            'overall_rate': min(overall_rate, 100),
            'total_required': total_required,
            'total_completed': total_completed
        }
    
    total_required = requirements.get('total_required', 10)
    
    total_completed = 0
    for category, stats in category_stats.items():
        completed = stats.get('completed', 0) + stats.get('assessed', 0)
        total_completed += completed
    
    overall_rate = (total_completed / total_required * 100) if total_required > 0 else 0
    overall_rate = min(overall_rate, 100)
    
    return {
        'overall_rate': overall_rate,
        'total_required': total_required,
        'total_completed': total_completed
    }

def identify_capability_weaknesses(person_name, capability_scores, threshold=70):
    """识别个人能力缺陷"""
    weaknesses = []
    
    if person_name not in capability_scores:
        return weaknesses
    
    person_scores = capability_scores[person_name]
    
    for category, score in person_scores['专业能力'].items():
        if score < threshold:
            weaknesses.append({
                'type': '专业能力',
                'category': category,
                'score': score,
                'threshold': threshold,
                'gap': threshold - score
            })
    
    for category, score in person_scores['核心能力'].items():
        short_name = category
        if '客户价值的挖掘' in category:
            short_name = '价值挖掘'
        elif '客户需求的捕获' in category:
            short_name = '需求捕获'
        elif '客户关系的维护' in category:
            short_name = '关系维护'
        elif '投诉抱怨' in category:
            short_name = '投诉化解'
        elif '资源统筹' in category:
            short_name = '资源统筹'
        
        weaknesses.append({
            'type': '核心能力',
            'category': short_name,
            'full_name': category,
            'score': score,
            'threshold': threshold,
            'gap': threshold - score
        })
    
    weaknesses.sort(key=lambda x: x['gap'], reverse=True)
    
    return weaknesses

def process_data(data_dict):
    """处理所有数据"""
    
    position_requirements = {}
    training_req_df = data_dict['岗位培训要求']
    
    for _, row in training_req_df.iterrows():
        if pd.notna(row.get('岗位')):
            position = str(row['岗位']).strip()
            total_required = row.get('培训课程总数', 10)
            try:
                total_required = int(total_required)
            except:
                total_required = 10
            
            position_requirements[position] = {
                'total_required': total_required
            }
    
    personnel_list = []
    personnel_info = {}
    skill_tags_all = []
    all_cases = []
    
    personnel_cols = data_dict['人员信息'].columns.tolist()
    
    for _, row in data_dict['人员信息'].iterrows():
        if pd.notna(row['姓名']):
            name = str(row['姓名']).strip()
            personnel_list.append(name)
            
            skill_tags = []
            if '能力标签' in personnel_cols:
                tags_value = row['能力标签']
                if pd.notna(tags_value):
                    skill_tags = parse_skill_tags(tags_value)
                    skill_tags_all.extend(skill_tags)
            
            classic_cases = []
            if '经典案例' in personnel_cols:
                case_value = row['经典案例']
                if pd.notna(case_value):
                    classic_cases = parse_classic_cases(case_value)
                    all_cases.extend([(name, case) for case in classic_cases])
            
            info_dict = {}
            for col in ['工号', '岗位', '岗位年限', '职称', '对接客户', '入职日期', '联系方式']:
                if col in personnel_cols:
                    info_dict[col] = str(row[col]) if pd.notna(row[col]) else ''
                else:
                    info_dict[col] = ''
            
            info_dict['能力标签'] = skill_tags
            info_dict['经典案例'] = classic_cases
            personnel_info[name] = info_dict
    
    tag_frequency = {}
    for tag in skill_tags_all:
        tag_frequency[tag] = tag_frequency.get(tag, 0) + 1
    
    training_stats = {}
    course_categories = {}
    
    for person in personnel_list:
        training_stats[person] = {
            'completed': 0,
            'assessed': 0,
            'not_completed': 0,
            'total': 0,
            'completion_rate': 0,
            'category_stats': {}
        }
    
    for _, row in data_dict['培训课程'].iterrows():
        if '培训大类' in row.index:
            category = str(row['培训大类']) if pd.notna(row['培训大类']) else '其他'
        else:
            category = '其他'
            
        if category not in course_categories:
            course_categories[category] = 0
        course_categories[category] += 1
        
        for person in personnel_list:
            if person in row.index:
                status = str(row[person]) if pd.notna(row[person]) else ''
                training_stats[person]['total'] += 1
                
                if category not in training_stats[person]['category_stats']:
                    training_stats[person]['category_stats'][category] = {'total': 0, 'completed': 0, 'assessed': 0}
                
                training_stats[person]['category_stats'][category]['total'] += 1
                
                if '已完成培训' in status or '已完成' in status:
                    training_stats[person]['completed'] += 1
                    training_stats[person]['category_stats'][category]['completed'] += 1
                elif '经评估具备该能力' in status or '经评估' in status:
                    training_stats[person]['assessed'] += 1
                    training_stats[person]['category_stats'][category]['assessed'] += 1
                elif '未完成培训' in status or '未完成' in status:
                    training_stats[person]['not_completed'] += 1
    
    for person in personnel_list:
        total = training_stats[person]['total']
        if total > 0:
            effective_completed = training_stats[person]['completed'] + training_stats[person]['assessed']
            training_stats[person]['completion_rate'] = round((effective_completed / total) * 100, 1)
    
    capability_scores = {}
    for person in personnel_list:
        capability_scores[person] = {
            '专业能力': {},
            '核心能力': {}
        }
    
    for _, row in data_dict['能力评估'].iterrows():
        if '姓名' in row.index and pd.notna(row['姓名']):
            person = str(row['姓名']).strip()
            if person in personnel_list:
                prof_categories = ['通识类', '商务类', '业务类', '产品类', '管理类']
                for category in prof_categories:
                    if category in row.index and pd.notna(row[category]):
                        try:
                            score = float(row[category])
                            capability_scores[person]['专业能力'][category] = min(100, max(0, score))
                        except:
                            capability_scores[person]['专业能力'][category] = 0
                
                core_mapping = {
                    '客户价值挖掘': '客户价值的挖掘与管理能力',
                    '客户需求捕获': '客户需求的捕获与处理能力',
                    '客户关系维护': '客户关系的维护与商务谈判能力',
                    '投诉抱怨化解': '客户潜在的投诉抱怨识别与化解能力',
                    '资源统筹能力': '内部各业务接口及资源统筹能力'
                }
                
                for short_name, full_name in core_mapping.items():
                    if short_name in row.index and pd.notna(row[short_name]):
                        try:
                            score = float(row[short_name])
                            capability_scores[person]['核心能力'][full_name] = min(100, max(0, score))
                        except:
                            capability_scores[person]['核心能力'][full_name] = 0
    
    position_training_completion = {}
    for person in personnel_list:
        position = personnel_info[person].get('岗位', '其他')
        category_stats = training_stats[person]['category_stats']
        position_training_completion[person] = calculate_training_completion_by_position(
            person, position, category_stats, position_requirements
        )
    
    return personnel_list, personnel_info, training_stats, capability_scores, course_categories, tag_frequency, all_cases, position_training_completion, position_requirements

def create_radar_chart(person_name, capability_scores):
    """创建能力雷达图"""
    if person_name not in capability_scores:
        return None
    
    person_scores = capability_scores[person_name]
    
    fig = make_subplots(
        rows=1, cols=2,
        specs=[[{'type': 'polar'}, {'type': 'polar'}]],
        subplot_titles=('五大专业能力评估', '五大核心能力评估'),
        horizontal_spacing=0.15
    )
    
    prof_categories = ['通识类', '商务类', '业务类', '产品类', '管理类']
    prof_values = [person_scores['专业能力'].get(cat, 0) for cat in prof_categories]
    
    core_categories_short = ['价值挖掘', '需求捕获', '关系维护', '投诉化解', '资源统筹']
    core_full_categories = list(person_scores['核心能力'].keys())
    core_values = [person_scores['核心能力'].get(cat, 0) for cat in core_full_categories]
    
    if not core_values:
        core_values = [0] * 5
    
    fig.add_trace(
        go.Scatterpolar(
            r=prof_values,
            theta=prof_categories,
            fill='toself',
            name='专业能力',
            line_color='#4F46E5',
            fillcolor='rgba(79, 70, 229, 0.3)',
            hovertemplate='<b>%{theta}</b><br>得分: %{r:.1f}分<extra></extra>'
        ),
        row=1, col=1
    )
    
    fig.add_trace(
        go.Scatterpolar(
            r=core_values,
            theta=core_categories_short,
            fill='toself',
            name='核心能力',
            line_color='#10B981',
            fillcolor='rgba(16, 185, 129, 0.3)',
            hovertemplate='<b>%{theta}</b><br>得分: %{r:.1f}分<extra></extra>'
        ),
        row=1, col=2
    )
    
    fig.update_layout(
        polar=dict(
            radialaxis=dict(
                visible=True,
                range=[0, 100],
                tickvals=[0, 25, 50, 75, 100],
                ticktext=['0', '25', '50', '75', '100'],
                tickfont=dict(size=10)
            ),
            bgcolor='rgba(0,0,0,0)'
        ),
        polar2=dict(
            radialaxis=dict(
                visible=True,
                range=[0, 100],
                tickvals=[0, 25, 50, 75, 100],
                ticktext=['0', '25', '50', '75', '100'],
                tickfont=dict(size=10)
            ),
            bgcolor='rgba(0,0,0,0)'
        ),
        showlegend=False,
        height=500,
        title_text=f"{person_name} - 能力评估雷达图",
        title_font=dict(size=20, color='#1E3A8A'),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    return fig

def create_comparison_radar_chart(person1_name, person2_name, capability_scores):
    """创建对比雷达图"""
    if person1_name not in capability_scores or person2_name not in capability_scores:
        return None
    
    person1_scores = capability_scores[person1_name]
    person2_scores = capability_scores[person2_name]
    
    fig = make_subplots(
        rows=1, cols=2,
        specs=[[{'type': 'polar'}, {'type': 'polar'}]],
        subplot_titles=(f'{person1_name} vs {person2_name} - 专业能力', f'{person1_name} vs {person2_name} - 核心能力'),
        horizontal_spacing=0.15
    )
    
    prof_categories = ['通识类', '商务类', '业务类', '产品类', '管理类']
    person1_prof_values = [person1_scores['专业能力'].get(cat, 0) for cat in prof_categories]
    person2_prof_values = [person2_scores['专业能力'].get(cat, 0) for cat in prof_categories]
    
    core_categories_short = ['价值挖掘', '需求捕获', '关系维护', '投诉化解', '资源统筹']
    core_full_categories = list(person1_scores['核心能力'].keys())
    person1_core_values = [person1_scores['核心能力'].get(cat, 0) for cat in core_full_categories]
    person2_core_values = [person2_scores['核心能力'].get(cat, 0) for cat in core_full_categories]
    
    fig.add_trace(
        go.Scatterpolar(
            r=person1_prof_values,
            theta=prof_categories,
            fill='toself',
            name=person1_name,
            line_color='#4F46E5',
            fillcolor='rgba(79, 70, 229, 0.3)',
            hovertemplate='<b>%{theta}</b><br>' + person1_name + ': %{r:.1f}分<extra></extra>'
        ),
        row=1, col=1
    )
    
    fig.add_trace(
        go.Scatterpolar(
            r=person2_prof_values,
            theta=prof_categories,
            fill='toself',
            name=person2_name,
            line_color='#10B981',
            fillcolor='rgba(16, 185, 129, 0.3)',
            hovertemplate='<b>%{theta}</b><br>' + person2_name + ': %{r:.1f}分<extra></extra>'
        ),
        row=1, col=1
    )
    
    fig.add_trace(
        go.Scatterpolar(
            r=person1_core_values,
            theta=core_categories_short,
            fill='toself',
            name=person1_name,
            line_color='#4F46E5',
            fillcolor='rgba(79, 70, 229, 0.3)',
            hovertemplate='<b>%{theta}</b><br>' + person1_name + ': %{r:.1f}分<extra></extra>',
            showlegend=False
        ),
        row=1, col=2
    )
    
    fig.add_trace(
        go.Scatterpolar(
            r=person2_core_values,
            theta=core_categories_short,
            fill='toself',
            name=person2_name,
            line_color='#10B981',
            fillcolor='rgba(16, 185, 129, 0.3)',
            hovertemplate='<b>%{theta}</b><br>' + person2_name + ': %{r:.1f}分<extra></extra>',
            showlegend=False
        ),
        row=1, col=2
    )
    
    fig.update_layout(
        polar=dict(
            radialaxis=dict(
                visible=True,
                range=[0, 100],
                tickvals=[0, 25, 50, 75, 100],
                ticktext=['0', '25', '50', '75', '100'],
                tickfont=dict(size=10)
            ),
            bgcolor='rgba(0,0,0,0)'
        ),
        polar2=dict(
            radialaxis=dict(
                visible=True,
                range=[0, 100],
                tickvals=[0, 25, 50, 75, 100],
                ticktext=['0', '25', '50', '75', '100'],
                tickfont=dict(size=10)
            ),
            bgcolor='rgba(0,0,0,0)'
        ),
        showlegend=True,
        legend=dict(
            yanchor="top",
            y=0.99,
            xanchor="left",
            x=1.02
        ),
        height=500,
        title_text="能力对比雷达图",
        title_font=dict(size=20, color='#1E3A8A'),
        paper_bgcolor='rgba(0,0,0,0)',
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    return fig

def display_skill_tags(tags):
    """显示能力标签"""
    if not tags:
        st.info("暂无能力标签")
        return
    
    html_tags = '<div class="tag-container">'
    for tag in tags:
        category = get_tag_category(tag)
        html_tags += f'<span class="skill-tag {category}">🏷️ {tag}</span>'
    html_tags += '</div>'
    
    st.markdown(html_tags, unsafe_allow_html=True)

def display_classic_cases(cases):
    """显示经典案例"""
    if not cases:
        st.info("暂无经典案例记录")
        return
    
    for i, case in enumerate(cases):
        icon = get_case_icon(case.get('type', '项目案例'))
        
        meta_items = []
        
        if case.get('year'):
            meta_items.append(f'<span>📅 {case["year"]}</span>')
        
        if case.get('role'):
            meta_items.append(f'<span>👤 参与角色: {case["role"]}</span>')
        
        if case.get('achievement'):
            meta_items.append(f'<span>🏆 成果: {case["achievement"]}</span>')
        
        if meta_items:
            meta_html = f'<div class="case-meta">{"".join(meta_items)}</div>'
        else:
            meta_html = ''
        
        case_html = f"""
        <div class="case-card">
            <div class="case-title">
                {icon} {case.get('title', '经典案例')}
                {f'<span class="case-tag" style="margin-left: auto;">{case.get("type", "项目案例")}</span>' if case.get('type') else ''}
            </div>
            <div class="case-content">
                {case.get('content', '')}
            </div>
            {meta_html}
        </div>
        """
        
        st.markdown(case_html, unsafe_allow_html=True)

def display_capability_weaknesses(weaknesses):
    """显示能力缺陷"""
    if not weaknesses:
        st.success("✅ 所有能力均达标！")
        return
    
    st.markdown("#### ⚠️ 待提升能力")
    
    for weakness in weaknesses[:5]:
        gap_percentage = (weakness['gap'] / weakness['threshold'] * 100)
        if gap_percentage >= 30:
            level = "严重不足"
        elif gap_percentage >= 15:
            level = "明显不足"
        else:
            level = "轻微不足"
        
        st.markdown(f"""
        <span class="weakness-tag">
            {weakness['category']}: {weakness['score']:.1f}分 
            (低于阈值{weakness['threshold']}分 {weakness['gap']:.1f}分 - {level})
        </span>
        """, unsafe_allow_html=True)
    
    if len(weaknesses) > 5:
        st.caption(f"... 还有{len(weaknesses) - 5}项能力待提升")

def generate_warnings(personnel_list, personnel_info, capability_scores, warning_thresholds):
    """生成预警信息"""
    warnings = []
    
    for person in personnel_list:
        if person in personnel_info and person in capability_scores:
            position = personnel_info[person].get('岗位', '其他')
            threshold = warning_thresholds.get(position, warning_thresholds.get('其他', 70))
            
            weaknesses = identify_capability_weaknesses(person, capability_scores, threshold)
            
            if weaknesses:
                severe_weaknesses = [w for w in weaknesses if w['gap'] >= threshold * 0.3]
                moderate_weaknesses = [w for w in weaknesses if threshold * 0.15 <= w['gap'] < threshold * 0.3]
                mild_weaknesses = [w for w in weaknesses if w['gap'] < threshold * 0.15]
                
                for weakness in severe_weaknesses[:3]:
                    warnings.append({
                        'type': '能力缺陷',
                        '级别': '🔴 严重',
                        '级别_class': 'high',
                        '内容': f"{person} ({position}) 的 {weakness['category']} 能力得分仅 {weakness['score']:.1f}分，低于阈值 {threshold}分 {weakness['gap']:.1f}分",
                        '建议': f"建议立即安排 {weakness['category']} 专项培训，制定提升计划",
                        '岗位': position,
                        '分数': weakness['score'],
                        'category': weakness['category']
                    })
                
                if severe_weaknesses and warnings:
                    continue
                
                for weakness in moderate_weaknesses[:2]:
                    warnings.append({
                        'type': '能力缺陷',
                        '级别': '🟡 一般',
                        '级别_class': 'medium',
                        '内容': f"{person} ({position}) 的 {weakness['category']} 能力得分 {weakness['score']:.1f}分，需关注提升",
                        '建议': f"建议参加 {weakness['category']} 相关培训课程",
                        '岗位': position,
                        '分数': weakness['score'],
                        'category': weakness['category']
                    })
                
                if not warnings and mild_weaknesses:
                    for weakness in mild_weaknesses[:1]:
                        warnings.append({
                            'type': '能力缺陷',
                            '级别': '🔵 轻微',
                            '级别_class': 'low',
                            '内容': f"{person} ({position}) 的 {weakness['category']} 能力略有不足，得分 {weakness['score']:.1f}分",
                            '建议': f"建议通过自我学习提升 {weakness['category']} 能力",
                            '岗位': position,
                            '分数': weakness['score'],
                            'category': weakness['category']
                        })
    
    warning_order = {'high': 0, 'medium': 1, 'low': 2}
    warnings.sort(key=lambda x: warning_order.get(x.get('级别_class', 'low'), 3))
    
    return warnings

def create_demo_data():
    """创建演示数据"""
    demo_personnel = pd.DataFrame({
        '姓名': ['张三', '李四', '王五', '赵六', '钱七', '孙八'],
        '工号': ['001', '002', '003', '004', '005', '006'],
        '岗位': ['高级客户经理', '客户经理', '客户经理助理', '客户经理', '高级客户经理', '客户经理助理'],
        '岗位年限': ['5', '3', '1', '2', '6', '1.5'],
        '职称': ['高级工程师', '工程师', '助理工程师', '工程师', '高级工程师', '助理工程师'],
        '对接客户': ['国航、东航', '南航', '海航', '川航', '深航、厦航', '山航'],
        '入职日期': ['2018-01-01', '2020-03-15', '2022-06-01', '2021-08-20', '2017-05-10', '2022-11-01'],
        '联系方式': ['13800138001', '13800138002', '13800138003', '13800138004', '13800138005', '13800138006'],
        '能力标签': [
            '英语专八,沟通能力强,项目管理,数据分析',
            '日语N1,客户关系维护,谈判技巧',
            '团队协作,学习能力强',
            '英语六级,技术理解,产品知识',
            '英语专八,PMP认证,团队管理,商务谈判',
            '学习能力强,沟通能力,基础产品知识'
        ],
        '经典案例': [
            '国航EIS系统保障|成功为国航提供EIS系统保障支持，确保系统全年无故障运行|2023|技术支持|项目经理|客户满意度提升20%',
            '南航客户培训项目|主导南航客户培训项目，培训100+名客户经理|2022|培训指导|项目负责人|培训通过率95%',
            '海航需求调研|参与海航需求调研项目，收集整理200+条客户需求|2023|客户服务|参与成员|形成完整需求报告',
            '川航系统优化|提出川航系统优化建议，被采纳并实施|2022|创新成果|建议提出者|系统效率提升15%',
            '深航战略合作|主导深航战略合作项目，达成年度合作金额5000万|2023|项目案例|项目经理|合作金额创新高',
            '山航技术支持|为山航提供技术支持，解决重大技术问题|2023|技术支持|技术骨干|系统稳定性提升'
        ]
    })
    
    demo_training = pd.DataFrame({
        '培训大类': ['通识类', '通识类', '商务类', '业务类', '产品类', '管理类', '商务类', '业务类'],
        '培训课程': ['公司文化', '职业道德', '商务礼仪', '航空业务', '产品知识', '团队管理', '谈判技巧', '客户需求分析'],
        '张三': ['已完成培训', '经评估具备该能力', '已完成培训', '已完成培训', '未完成培训', '已完成培训', '已完成培训', '已完成培训'],
        '李四': ['已完成培训', '已完成培训', '经评估具备该能力', '未完成培训', '已完成培训', '未完成培训', '已完成培训', '已完成培训'],
        '王五': ['未完成培训', '已完成培训', '未完成培训', '未完成培训', '未完成培训', '未完成培训', '未完成培训', '已完成培训'],
        '赵六': ['已完成培训', '已完成培训', '已完成培训', '已完成培训', '已完成培训', '经评估具备该能力', '已完成培训', '已完成培训'],
        '钱七': ['已完成培训', '已完成培训', '已完成培训', '已完成培训', '已完成培训', '已完成培训', '已完成培训', '已完成培训'],
        '孙八': ['已完成培训', '已完成培训', '未完成培训', '已完成培训', '未完成培训', '未完成培训', '未完成培训', '已完成培训']
    })
    
    demo_capability = pd.DataFrame({
        '姓名': ['张三', '李四', '王五', '赵六', '钱七', '孙八'],
        '通识类': [85, 78, 65, 82, 92, 70],
        '商务类': [90, 85, 70, 88, 95, 72],
        '业务类': [88, 82, 68, 85, 90, 68],
        '产品类': [92, 80, 72, 90, 94, 65],
        '管理类': [87, 75, 60, 83, 91, 62],
        '客户价值挖掘': [85, 80, 65, 82, 90, 68],
        '客户需求捕获': [88, 82, 70, 85, 92, 70],
        '客户关系维护': [90, 85, 68, 88, 93, 72],
        '投诉抱怨化解': [86, 78, 62, 80, 89, 65],
        '资源统筹能力': [89, 80, 65, 84, 91, 66]
    })
    
    demo_requirements = pd.DataFrame({
        '岗位': ['高级客户经理', '客户经理', '客户经理助理'],
        '培训课程总数': [14, 10, 7]
    })
    
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        demo_personnel.to_excel(writer, sheet_name='人员信息', index=False)
        demo_training.to_excel(writer, sheet_name='培训课程', index=False)
        demo_capability.to_excel(writer, sheet_name='能力评估', index=False)
        demo_requirements.to_excel(writer, sheet_name='岗位培训要求', index=False)
    
    output.seek(0)
    return output

def generate_team_report(personnel_list, personnel_info, training_stats, capability_scores, all_cases, warnings_list, position_training_completion):
    """生成团队报告"""
    report_content = f"""
# 客户经理团队分析报告
生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

## 一、团队概况
- 团队成员总数: {len(personnel_list)}人
- 经典案例总数: {len(all_cases)}个
- 当前预警数量: {len(warnings_list)}个

## 二、岗位分布
"""
    
    positions = {}
    for info in personnel_info.values():
        position = info.get('岗位', '未知')
        positions[position] = positions.get(position, 0) + 1
    
    for position, count in positions.items():
        report_content += f"- {position}: {count}人\n"
    
    report_content += f"""
## 三、培训完成情况（按岗位要求）
"""
    
    for person in personnel_list:
        position = personnel_info[person].get('岗位', '其他')
        completion = position_training_completion[person]
        report_content += f"- {person} ({position}): 完成率 {completion['overall_rate']:.1f}% ({completion['total_completed']}/{completion['total_required']})\n"
    
    if warnings_list:
        report_content += f"""
## 四、能力缺陷预警
当前共有 {len(warnings_list)} 个预警:
"""
        
        high_warnings = [w for w in warnings_list if w['级别_class'] == 'high']
        medium_warnings = [w for w in warnings_list if w['级别_class'] == 'medium']
        low_warnings = [w for w in warnings_list if w['级别_class'] == 'low']
        
        if high_warnings:
            report_content += "\n### 严重预警:\n"
            for warning in high_warnings[:5]:
                report_content += f"- {warning['内容']}\n"
        
        if medium_warnings:
            report_content += "\n### 一般预警:\n"
            for warning in medium_warnings[:5]:
                report_content += f"- {warning['内容']}\n"
        
        if low_warnings:
            report_content += "\n### 轻微预警:\n"
            for warning in low_warnings[:5]:
                report_content += f"- {warning['内容']}\n"
    
    report_content += f"""
## 五、能力分析
"""
    
    if capability_scores:
        prof_categories = ['通识类', '商务类', '业务类', '产品类', '管理类']
        for category in prof_categories:
            scores = []
            for person in personnel_list:
                if person in capability_scores and category in capability_scores[person]['专业能力']:
                    scores.append(capability_scores[person]['专业能力'][category])
            if scores:
                avg_score = np.mean(scores)
                report_content += f"- {category}平均分: {avg_score:.1f}\n"
    
    report_content += f"""
## 六、改进建议

### 能力提升方面:
"""
    
    weakness_stats = {}
    for warning in warnings_list:
        if warning['type'] == '能力缺陷':
            content = warning['内容']
            for category in ['通识类', '商务类', '业务类', '产品类', '管理类', 
                           '价值挖掘', '需求捕获', '关系维护', '投诉化解', '资源统筹']:
                if category in content:
                    weakness_stats[category] = weakness_stats.get(category, 0) + 1
                    break
    
    if weakness_stats:
        report_content += "\n团队共性能力短板:\n"
        sorted_weaknesses = sorted(weakness_stats.items(), key=lambda x: x[1], reverse=True)
        for category, count in sorted_weaknesses[:3]:
            report_content += f"- {category}: {count}人需要提升\n"
        
        report_content += "\n建议组织相关专题培训:\n"
        for category, _ in sorted_weaknesses[:2]:
            report_content += f"- {category}强化培训\n"
    
    report_content += """
### 培训管理方面:
1. 根据岗位要求动态调整培训计划
2. 建立能力提升跟踪机制
3. 定期评估培训效果

---
报告结束
"""
    
    return report_content

def get_default_photo():
    """获取默认头像"""
    img = Image.new('RGB', (200, 250), color=(59, 130, 246))
    return img

def handle_photo_upload(person_name):
    """处理照片上传"""
    col1, col2 = st.columns([2, 1])
    
    with col1:
        uploaded_photo = st.file_uploader(f"为 {person_name} 上传证件照", 
                                         type=['jpg', 'jpeg', 'png'],
                                         key=f"photo_{person_name}")
        
        if uploaded_photo is not None:
            try:
                image = Image.open(uploaded_photo)
                image = image.resize((200, 250))
                st.session_state.personnel_photos[person_name] = image
                st.success("✅ 照片上传成功！")
            except Exception as e:
                st.error(f"❌ 图片处理失败: {str(e)}")
    
    with col2:
        if person_name in st.session_state.personnel_photos:
            st.image(st.session_state.personnel_photos[person_name], 
                    caption=f"{person_name} 的证件照", 
                    width=150)
        else:
            st.info("尚未上传照片")
            default_img = get_default_photo()
            st.image(default_img, caption="默认头像", width=150)

def generate_team_analysis_report(personnel_list, personnel_info, training_stats, capability_scores, all_cases, position_training_completion, warnings_list, warning_thresholds, team_threshold):
    """生成增强的团队分析报告"""
    
    total_people = len(personnel_list)
    total_cases = len(all_cases)
    
    total_skill_tags = sum([len(info['能力标签']) for info in personnel_info.values()])
    unique_skill_tags = set()
    for info in personnel_info.values():
        unique_skill_tags.update(info['能力标签'])
    
    completion_groups_with_names = {
        '优秀 (≥90%)': [],
        '良好 (70-89%)': [],
        '待提升 (50-69%)': [],
        '需关注 (<50%)': []
    }
    
    position_completion_rates = []
    
    for person in personnel_list:
        if person in position_training_completion:
            completion = position_training_completion[person]
            position_completion_rates.append(completion['overall_rate'])
            
            rate = completion['overall_rate']
            person_info = {
                '姓名': person,
                '岗位': personnel_info[person].get('岗位', '其他'),
                '完成率': rate,
                '进度': f"{completion['total_completed']}/{completion['total_required']}"
            }
            
            if rate >= 90:
                completion_groups_with_names['优秀 (≥90%)'].append(person_info)
            elif rate >= 70:
                completion_groups_with_names['良好 (70-89%)'].append(person_info)
            elif rate >= 50:
                completion_groups_with_names['待提升 (50-69%)'].append(person_info)
            else:
                completion_groups_with_names['需关注 (<50%)'].append(person_info)
    
    avg_completion_rate = np.mean(position_completion_rates) if position_completion_rates else 0
    
    avg_prof_scores = {}
    avg_core_scores = {}
    
    prof_categories = ['通识类', '商务类', '业务类', '产品类', '管理类']
    core_categories = ['客户价值的挖掘与管理能力', '客户需求的捕获与处理能力', 
                      '客户关系的维护与商务谈判能力', '客户潜在的投诉抱怨识别与化解能力',
                      '内部各业务接口及资源统筹能力']
    
    for category in prof_categories:
        scores = []
        for person in personnel_list:
            if person in capability_scores and category in capability_scores[person]['专业能力']:
                scores.append(capability_scores[person]['专业能力'][category])
        if scores:
            avg_prof_scores[category] = np.mean(scores)
    
    for category in core_categories:
        scores = []
        for person in personnel_list:
            if person in capability_scores and category in capability_scores[person]['核心能力']:
                scores.append(capability_scores[person]['核心能力'][category])
        if scores:
            avg_core_scores[category] = np.mean(scores)
    
    positions = {}
    titles = {}
    for info in personnel_info.values():
        position = info.get('岗位', '未知')
        if position:
            positions[position] = positions.get(position, 0) + 1
        
        title = info.get('职称', '未知')
        if title:
            titles[title] = titles.get(title, 0) + 1
    
    skill_tag_counts = {}
    for info in personnel_info.values():
        for tag in info['能力标签']:
            skill_tag_counts[tag] = skill_tag_counts.get(tag, 0) + 1
    
    case_years = {}
    for _, case in all_cases:
        year = case.get('year', '')
        if year:
            case_years[year] = case_years.get(year, 0) + 1
    
    team_weaknesses = {}
    for person in personnel_list:
        if person in capability_scores:
            for category, score in capability_scores[person]['专业能力'].items():
                if score < team_threshold:
                    if category not in team_weaknesses:
                        team_weaknesses[category] = {
                            'count': 0,
                            'persons': [],
                            'avg_score': 0,
                            'total_score': 0
                        }
                    team_weaknesses[category]['count'] += 1
                    team_weaknesses[category]['persons'].append(person)
                    team_weaknesses[category]['total_score'] += score
            
            for category, score in capability_scores[person]['核心能力'].items():
                short_name = category
                if '客户价值的挖掘' in category:
                    short_name = '价值挖掘'
                elif '客户需求的捕获' in category:
                    short_name = '需求捕获'
                elif '客户关系的维护' in category:
                    short_name = '关系维护'
                elif '投诉抱怨' in category:
                    short_name = '投诉化解'
                elif '资源统筹' in category:
                    short_name = '资源统筹'
                
                if score < team_threshold:
                    if short_name not in team_weaknesses:
                        team_weaknesses[short_name] = {
                            'count': 0,
                            'persons': [],
                            'avg_score': 0,
                            'total_score': 0
                        }
                    team_weaknesses[short_name]['count'] += 1
                    team_weaknesses[short_name]['persons'].append(person)
                    team_weaknesses[short_name]['total_score'] += score
    
    for category, data in team_weaknesses.items():
        data['avg_score'] = data['total_score'] / data['count'] if data['count'] > 0 else 0
    
    sorted_team_weaknesses = sorted(team_weaknesses.items(), key=lambda x: x[1]['count'], reverse=True)
    
    return {
        'total_people': total_people,
        'total_cases': total_cases,
        'avg_completion_rate': avg_completion_rate,
        'total_skill_tags': total_skill_tags,
        'unique_skill_tags_count': len(unique_skill_tags),
        'completion_groups_with_names': completion_groups_with_names,
        'avg_prof_scores': avg_prof_scores,
        'avg_core_scores': avg_core_scores,
        'positions': positions,
        'titles': titles,
        'skill_tag_counts': skill_tag_counts,
        'case_years': case_years,
        'unique_skill_tags': list(unique_skill_tags),
        'team_weaknesses': sorted_team_weaknesses,
        'warnings_list': warnings_list
    }

def create_training_completion_chart(completion_groups, completion_groups_with_names):
    """创建培训完成情况分布图（带人名悬停）"""
    categories = list(completion_groups.keys())
    values = list(completion_groups.values())
    
    colors = {
        '优秀 (≥90%)': '#10B981',
        '良好 (70-89%)': '#3B82F6',
        '待提升 (50-69%)': '#F59E0B',
        '需关注 (<50%)': '#EF4444'
    }
    
    marker_colors = [colors.get(cat, '#6B7280') for cat in categories]
    
    hover_texts = []
    for category in categories:
        persons = completion_groups_with_names.get(category, [])
        if persons:
            text = f"<b>{category}</b><br>人数: {len(persons)}<br><br><b>人员名单:</b><br>"
            for person in persons:
                text += f"👤 {person['姓名']} ({person['岗位']}) - {person['完成率']:.1f}% ({person['进度']})<br>"
            hover_texts.append(text)
        else:
            hover_texts.append(f"<b>{category}</b><br>人数: 0")
    
    fig = go.Figure(data=[
        go.Bar(
            x=categories,
            y=values,
            text=values,
            textposition='auto',
            marker_color=marker_colors,
            hovertemplate='%{customdata}<extra></extra>',
            customdata=hover_texts,
            width=[0.3]
        )
    ])
    
    fig.update_layout(
        title='培训完成率分布情况（按岗位要求）',
        xaxis_title='完成率分组',
        yaxis_title='人数',
        height=400,
        plot_bgcolor='rgba(0,0,0,0)',
        showlegend=False,
        bargap=0.5
    )
    
    return fig

def display_team_weaknesses_with_threshold(team_weaknesses, team_threshold):
    """显示团队能力短板（使用统一的阈值）"""
    if not team_weaknesses:
        st.success("✅ 团队整体能力表现良好，无明显短板！")
        return
    
    st.markdown("#### 🎯 团队能力短板分析")
    st.markdown(f"**当前分析阈值: {team_threshold}分**")
    
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("总短板数量", len(team_weaknesses))
    with col2:
        total_affected = sum([data[1]['count'] for data in team_weaknesses])
        st.metric("受影响总人次", total_affected)
    with col3:
        avg_gap = np.mean([(team_threshold - data[1]['avg_score']) for data in team_weaknesses]) if team_weaknesses else 0
        st.metric("平均差距", f"{avg_gap:.1f}分")
    
    for category, data in team_weaknesses[:5]:
        if data['count'] >= 3:
            severity = "严重"
            color = "#DC2626"
        elif data['count'] >= 2:
            severity = "明显"
            color = "#F59E0B"
        else:
            severity = "轻微"
            color = "#3B82F6"
        
        gap = team_threshold - data['avg_score']
        
        st.markdown(f"""
        <div class="team-weakness-card">
            <div class="team-weakness-title" style="color: {color}">
                {category} - {severity}短板 ({data['count']}人需要提升)
            </div>
            <div>
                <span style="font-weight: 500;">平均得分: {data['avg_score']:.1f}分</span>
                <span style="margin-left: 15px; color: #666;">(低于阈值{team_threshold}分 {gap:.1f}分)</span>
            </div>
            <div style="margin-top: 10px;">
        """, unsafe_allow_html=True)
        
        for person in data['persons'][:5]:
            st.markdown(f'<span class="team-weakness-item">👤 {person}</span>', unsafe_allow_html=True)
        
        if len(data['persons']) > 5:
            st.caption(f"... 还有{len(data['persons']) - 5}人")
        
        st.markdown('</div></div>', unsafe_allow_html=True)

def create_position_distribution_chart(positions):
    """创建岗位分布图"""
    labels = list(positions.keys())
    values = list(positions.values())
    
    fig = go.Figure(data=[
        go.Pie(
            labels=labels,
            values=values,
            hole=0.4,
            marker_colors=px.colors.qualitative.Set3,
            textinfo='label+percent',
            hovertemplate='<b>%{label}</b><br>人数: %{value} (%{percent})<extra></extra>'
        )
    ])
    
    fig.update_layout(
        title='岗位分布情况',
        height=400,
        showlegend=False
    )
    
    return fig

def create_title_distribution_chart(titles):
    """创建职称分布图"""
    labels = list(titles.keys())
    values = list(titles.values())
    
    fig = go.Figure(data=[
        go.Pie(
            labels=labels,
            values=values,
            hole=0.4,
            marker_colors=px.colors.qualitative.Pastel,
            textinfo='label+percent',
            hovertemplate='<b>%{label}</b><br>人数: %{value} (%{percent})<extra></extra>'
        )
    ])
    
    fig.update_layout(
        title='职称分布情况',
        height=400,
        showlegend=False
    )
    
    return fig

def create_capability_comparison_chart(avg_prof_scores, avg_core_scores, team_threshold):
    """创建能力对比图（带阈值线）"""
    prof_categories = list(avg_prof_scores.keys())
    prof_values = list(avg_prof_scores.values())
    
    core_categories_short = []
    core_full_names = list(avg_core_scores.keys())
    for full_name in core_full_names:
        if '价值挖掘' in full_name:
            core_categories_short.append('价值挖掘')
        elif '需求捕获' in full_name:
            core_categories_short.append('需求捕获')
        elif '关系维护' in full_name:
            core_categories_short.append('关系维护')
        elif '投诉化解' in full_name:
            core_categories_short.append('投诉化解')
        elif '资源统筹' in full_name:
            core_categories_short.append('资源统筹')
        else:
            core_categories_short.append(full_name[:4])
    
    core_values = list(avg_core_scores.values())
    
    fig = make_subplots(
        rows=1, cols=2,
        subplot_titles=('专业能力平均得分', '核心能力平均得分')
    )
    
    fig.add_trace(
        go.Bar(
            x=prof_categories,
            y=prof_values,
            name='专业能力',
            marker_color='#4F46E5',
            text=[f'{v:.1f}' for v in prof_values],
            textposition='auto',
            hovertemplate='<b>%{x}</b><br>平均分: %{y:.1f}<extra></extra>'
        ),
        row=1, col=1
    )
    
    fig.add_trace(
        go.Bar(
            x=core_categories_short,
            y=core_values,
            name='核心能力',
            marker_color='#10B981',
            text=[f'{v:.1f}' for v in core_values],
            textposition='auto',
            hovertemplate='<b>%{x}</b><br>平均分: %{y:.1f}<extra></extra>'
        ),
        row=1, col=2
    )
    
    fig.add_hline(y=team_threshold, line_dash="dash", line_color="#EF4444", 
                  annotation_text=f"阈值: {team_threshold}分", annotation_position="top left",
                  row=1, col=1)
    fig.add_hline(y=team_threshold, line_dash="dash", line_color="#EF4444", 
                  annotation_text=f"阈值: {team_threshold}分", annotation_position="top left",
                  row=1, col=2)
    
    fig.update_layout(
        height=400,
        showlegend=False,
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    fig.update_yaxes(range=[0, 100], row=1, col=1, title='得分')
    fig.update_yaxes(range=[0, 100], row=1, col=2, title='得分')
    
    return fig

def create_skill_tag_bar_chart(skill_tag_counts, top_n=10):
    """创建技能标签条形图"""
    sorted_tags = sorted(skill_tag_counts.items(), key=lambda x: x[1], reverse=True)
    top_tags = sorted_tags[:top_n]
    
    if not top_tags:
        return None
    
    tags = [tag[0] for tag in top_tags]
    counts = [tag[1] for tag in top_tags]
    
    fig = go.Figure(data=[
        go.Bar(
            x=counts,
            y=tags,
            orientation='h',
            marker_color='#8B5CF6',
            text=counts,
            textposition='outside',
            hovertemplate='<b>%{y}</b><br>掌握人数: %{x}<extra></extra>'
        )
    ])
    
    fig.update_layout(
        title=f'Top {top_n} 热门能力标签',
        xaxis_title='掌握人数',
        yaxis_title='能力标签',
        height=400 + (top_n * 20),
        plot_bgcolor='rgba(0,0,0,0)'
    )
    
    return fig

def export_all_to_pdf(personnel_list, personnel_info, capability_scores, training_stats, 
                      position_training_completion, all_cases, warnings_list, team_report):
    """标准PDF导出（仅文本）"""
    if not REPORTLAB_AVAILABLE:
        return None
    
    pdf_buffer = io.BytesIO()
    doc = SimpleDocTemplate(pdf_buffer, pagesize=A4, rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72)
    
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'], fontSize=24, textColor=colors.HexColor('#1E3A8A'), alignment=TA_CENTER, spaceAfter=30)
    heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'], fontSize=18, textColor=colors.HexColor('#4F46E5'), spaceAfter=12, spaceBefore=20)
    normal_style = ParagraphStyle('CustomNormal', parent=styles['Normal'], fontSize=10, textColor=colors.HexColor('#4B5563'), spaceAfter=6)
    
    story = []
    story.append(Paragraph("客户经理培训档案系统", title_style))
    story.append(Paragraph(f"生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", normal_style))
    
    doc.build(story)
    pdf_buffer.seek(0)
    return pdf_buffer

def export_all_to_ppt(personnel_list, personnel_info, capability_scores, training_stats, 
                      position_training_completion, all_cases, warnings_list, team_report):
    """PPT导出"""
    if not PPTX_AVAILABLE:
        return None
    
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "客户经理培训档案系统"
    
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer

# ========== 主函数 ==========
def main():
    """主函数"""
    
    # 尝试从本地缓存加载数据（仅在 session_state 未加载时）
    if not st.session_state.data_loaded:
        cached_data, cached_name = load_cached_data()
        if cached_data is not None:
            st.session_state.analysis_data = cached_data
            st.session_state.uploaded_file_name = cached_name
            st.session_state.data_loaded = True
            # 可选：在侧边栏显示提示
            # st.sidebar.success(f"已从本地缓存加载数据（上次文件：{cached_name}）")
    
    st.markdown('<div class="main-header">✈️ 客户经理培训档案管理系统</div>', unsafe_allow_html=True)
    
    with st.sidebar:
        st.image("https://img.icons8.com/color/96/000000/airplane-front-view.png", width=80)
        st.markdown("### 📤 数据上传")
        
        if st.button("📥 下载演示数据模板", key="download_demo"):
            demo_data = create_demo_data()
            st.download_button(
                label="💾 点击下载演示数据",
                data=demo_data,
                file_name="客户经理培训档案_演示数据.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="download_demo_file"
            )
        
        # 根据数据加载状态显示不同 UI
        if not st.session_state.data_loaded:
            # 未加载数据：显示文件上传器
            uploaded_file = st.file_uploader("选择Excel文件", type=['xlsx', 'xls'], key="file_uploader")
            
            if uploaded_file is not None:
                with st.spinner('📊 正在分析数据...'):
                    data_dict = load_excel_file(uploaded_file)
                    if data_dict:
                        personnel_list, personnel_info, training_stats, capability_scores, course_categories, tag_frequency, all_cases, position_training_completion, position_requirements = process_data(data_dict)
                        
                        # 保存到 session_state
                        st.session_state.analysis_data = {
                            'data_dict': data_dict,
                            'personnel_list': personnel_list,
                            'personnel_info': personnel_info,
                            'training_stats': training_stats,
                            'capability_scores': capability_scores,
                            'course_categories': course_categories,
                            'tag_frequency': tag_frequency,
                            'all_cases': all_cases,
                            'position_training_completion': position_training_completion,
                            'position_requirements': position_requirements
                        }
                        st.session_state.uploaded_file_name = uploaded_file.name
                        st.session_state.data_loaded = True
                        
                        # 保存到本地缓存文件
                        save_cached_data(st.session_state.analysis_data, uploaded_file.name)
                        
                        st.success("✅ 文件上传成功！")
                        st.rerun()
        else:
            # 已加载数据：显示文件信息和更换按钮
            st.info(f"当前文件：{st.session_state.uploaded_file_name}")
            if st.button("🔄 更换文件", key="change_file"):
                # 删除缓存文件
                if os.path.exists(CACHE_FILE):
                    os.remove(CACHE_FILE)
                # 清除 session_state
                st.session_state.data_loaded = False
                st.session_state.uploaded_file_name = None
                st.session_state.analysis_data = None
                st.rerun()
        
        st.markdown("### ⚠️ 个人预警设置")
        st.markdown("#### (基于个人岗位)")
        
        st.markdown("**高级客户经理预警阈值:**")
        st.session_state.warning_thresholds['高级客户经理'] = st.slider(
            "高级客户经理", 0, 100, st.session_state.warning_thresholds['高级客户经理'],
            key="senior_threshold"
        )
        
        st.markdown("**客户经理预警阈值:**")
        st.session_state.warning_thresholds['客户经理'] = st.slider(
            "客户经理", 0, 100, st.session_state.warning_thresholds['客户经理'],
            key="manager_threshold"
        )
        
        st.markdown("**客户经理助理预警阈值:**")
        st.session_state.warning_thresholds['客户经理助理'] = st.slider(
            "客户经理助理", 0, 100, st.session_state.warning_thresholds['客户经理助理'],
            key="assistant_threshold"
        )
        
        st.markdown("**其他岗位预警阈值:**")
        st.session_state.warning_thresholds['其他'] = st.slider(
            "其他", 0, 100, st.session_state.warning_thresholds['其他'],
            key="other_threshold"
        )
        
        st.markdown("### 🎯 分析设置")
        show_advanced = st.checkbox("显示高级分析", value=True)
        enable_warnings = st.checkbox("启用智能预警", value=True)
        
        st.markdown("### 🎯 团队阈值设置")
        st.markdown("#### (独立于个人预警)")
        team_threshold = st.slider(
            "团队能力分析阈值", 0, 100, st.session_state.team_threshold,
            key="team_threshold_slider",
            help="用于团队能力短板分析的统一阈值"
        )
        st.session_state.team_threshold = team_threshold
        
        st.markdown("---")
        st.markdown("""
        <div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                    color: white; padding: 15px; border-radius: 10px;">
            <h4 style="margin: 0;">🚀 增强功能</h4>
            <p style="margin: 5px 0; font-size: 0.9rem;">• 柱状图悬停显示人名</p>
            <p style="margin: 5px 0; font-size: 0.9rem;">• 独立的团队阈值抓取</p>
            <p style="margin: 5px 0; font-size: 0.9rem;">• 阈值线可视化</p>
            <p style="margin: 5px 0; font-size: 0.9rem;">• 团队能力短板分析</p>
            <p style="margin: 5px 0; font-size: 0.9rem;">• 截图PDF导出（完美保留图表）</p>
        </div>
        """, unsafe_allow_html=True)
    
    if not st.session_state.data_loaded:
        st.markdown("""
        <div style="max-width: 900px; margin: 0 auto; padding: 20px;">
            <div style="text-align: center; margin-bottom: 3rem;">
                <h1 style="color: #1E3A8A;">✈️ 客户经理培训档案管理系统</h1>
                <p style="color: #6B7280; font-size: 1.2rem;">全面记录、智能分析、专业展示客户经理成长轨迹</p>
            </div>
            
            <div style="display: grid; grid-template-columns: repeat(3, 1fr); gap: 25px; margin-bottom: 3rem;">
                <div style="background: linear-gradient(135deg, #e0e7ff 0%, #c7d2fe 100%); 
                            padding: 25px; border-radius: 15px; text-align: center;">
                    <div style="font-size: 3rem; margin-bottom: 15px;">📸</div>
                    <h3 style="margin: 0 0 10px 0; color: #4F46E5;">证件照管理</h3>
                    <p style="color: #6B7280;">支持上传员工证件照，完善档案信息</p>
                </div>
                
                <div style="background: linear-gradient(135deg, #fce7f3 0%, #fbcfe8 100%); 
                            padding: 25px; border-radius: 15px; text-align: center;">
                    <div style="font-size: 3rem; margin-bottom: 15px;">📊</div>
                    <h3 style="margin: 0 0 10px 0; color: #BE185D;">柱状图悬停</h3>
                    <p style="color: #6B7280;">鼠标悬停显示详细人员名单</p>
                </div>
                
                <div style="background: linear-gradient(135deg, #dcfce7 0%, #bbf7d0 100%); 
                            padding: 25px; border-radius: 15px; text-align: center;">
                    <div style="font-size: 3rem; margin-bottom: 15px;">📋</div>
                    <h3 style="margin: 0 0 10px 0; color: #166534;">截图导出</h3>
                    <p style="color: #6B7280;">完美保留所有图表和样式</p>
                </div>
            </div>
            
            <div style="background: #f8fafc; padding: 30px; border-radius: 15px; margin-bottom: 2rem;">
                <h3 style="color: #1E3A8A; margin-top: 0;">📋 使用说明</h3>
                
                <h4>🎯 快速开始：</h4>
                <p>1. 点击左侧"下载演示数据模板"按钮获取模板文件</p>
                <p>2. 按照模板格式填写您的数据</p>
                <p>3. 上传Excel文件开始分析</p>
                
                <h4>📝 Excel模板要求：</h4>
                <p>• 必须包含四个工作表：<strong>人员信息</strong>、<strong>培训课程</strong>、<strong>能力评估</strong>、<strong>岗位培训要求</strong></p>
                <p>• 人员信息表需包含：姓名、工号、岗位、岗位年限、职称、对接客户等基本信息</p>
                <p>• 新增<strong>能力标签</strong>列：用逗号分隔多个标签，如"英语专八,沟通能力强,数据分析"</p>
                <p>• 新增<strong>经典案例</strong>列：使用竖线分隔，格式为"标题|内容|年份|类型|角色|成果"</p>
                <p>• <strong>岗位培训要求</strong>表需包含：岗位、培训课程总数两列</p>
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        st.markdown("### 🎬 演示数据预览")
        demo_data = create_demo_data()
        demo_data.seek(0)
        demo_dict = load_excel_file(demo_data)
        
        if demo_dict:
            with st.expander("👥 人员信息表（示例）"):
                st.dataframe(demo_dict['人员信息'], use_container_width=True)
            
            with st.expander("📚 培训课程表（示例）"):
                st.dataframe(demo_dict['培训课程'], use_container_width=True)
            
            with st.expander("📊 能力评估表（示例）"):
                st.dataframe(demo_dict['能力评估'], use_container_width=True)
            
            with st.expander("📋 岗位培训要求（示例）"):
                st.dataframe(demo_dict['岗位培训要求'], use_container_width=True)
    
    else:
        # 从 session_state 读取数据
        analysis = st.session_state.analysis_data
        data_dict = analysis['data_dict']
        personnel_list = analysis['personnel_list']
        personnel_info = analysis['personnel_info']
        training_stats = analysis['training_stats']
        capability_scores = analysis['capability_scores']
        course_categories = analysis['course_categories']
        tag_frequency = analysis['tag_frequency']
        all_cases = analysis['all_cases']
        position_training_completion = analysis['position_training_completion']
        position_requirements = analysis['position_requirements']
        
        # 根据当前阈值和选项动态计算预警和团队报告
        if enable_warnings:
            warnings_list = generate_warnings(
                personnel_list, personnel_info, capability_scores, st.session_state.warning_thresholds
            )
        else:
            warnings_list = []
        
        team_report = generate_team_analysis_report(
            personnel_list, personnel_info, training_stats, capability_scores, all_cases,
            position_training_completion, warnings_list, st.session_state.warning_thresholds,
            st.session_state.team_threshold
        )
        
        tab1, tab2, tab3, tab4, tab5 = st.tabs(["📈 数据概览", "👤 个人档案", "🔄 对比分析", "👥 团队分析", "📋 数据报告"])
        
        with tab1:
            st.markdown('<div class="section-title">📊 数据概览</div>', unsafe_allow_html=True)
            
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.markdown(f"""
                <div class="metric-card">
                    <div style="font-size: 0.9rem; color: #666;">总人数</div>
                    <div style="font-size: 2rem; font-weight: bold; color: #1E3A8A;">{len(personnel_list)}</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col2:
                total_cases = len(all_cases)
                st.markdown(f"""
                <div class="metric-card">
                    <div style="font-size: 0.9rem; color: #666;">经典案例数</div>
                    <div style="font-size: 2rem; font-weight: bold; color: #10B981;">{total_cases}</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col3:
                avg_rate = team_report['avg_completion_rate']
                st.markdown(f"""
                <div class="metric-card">
                    <div style="font-size: 0.9rem; color: #666;">平均培训完成度</div>
                    <div style="font-size: 2rem; font-weight: bold; color: #10B981;">{avg_rate:.1f}%</div>
                </div>
                """, unsafe_allow_html=True)
            
            with col4:
                warning_count = len(warnings_list) if enable_warnings and warnings_list else 0
                warning_color = "#EF4444" if warning_count > 0 else "#10B981"
                st.markdown(f"""
                <div class="metric-card">
                    <div style="font-size: 0.9rem; color: #666;">个人预警数量</div>
                    <div style="font-size: 2rem; font-weight: bold; color: {warning_color};">{warning_count}</div>
                </div>
                """, unsafe_allow_html=True)
            
            if enable_warnings and warnings_list:
                st.markdown('<div class="section-title">⚠️ 个人能力缺陷预警</div>', unsafe_allow_html=True)
                
                for warning in warnings_list:
                    warning_class = warning.get('级别_class', 'medium')
                    
                    st.markdown(f"""
                    <div class="warning-box {warning_class}">
                        <div style="display: flex; justify-content: space-between; align-items: center;">
                            <div>
                                <strong>{warning['级别']} - {warning['type']}</strong>
                                <p style="margin: 5px 0;">{warning['内容']}</p>
                            </div>
                        </div>
                        <p style="margin: 5px 0 0 0; font-size: 0.9rem; color: #666;">
                            💡 <strong>建议:</strong> {warning['建议']}
                        </p>
                    </div>
                    """, unsafe_allow_html=True)
            
            st.markdown('<div class="section-title">📚 课程分类分布</div>', unsafe_allow_html=True)
            
            if course_categories:
                fig = px.pie(
                    values=list(course_categories.values()),
                    names=list(course_categories.keys()),
                    hole=0.3,
                    color_discrete_sequence=px.colors.qualitative.Set3
                )
                fig.update_layout(height=400)
                st.plotly_chart(fig, use_container_width=True)
            
            if position_requirements:
                st.markdown('<div class="section-title">📋 岗位培训要求</div>', unsafe_allow_html=True)
                
                req_data = []
                for position, reqs in position_requirements.items():
                    req_data.append({
                        '岗位': position,
                        '培训课程总数': reqs['total_required']
                    })
                
                if req_data:
                    req_df = pd.DataFrame(req_data)
                    st.dataframe(req_df, use_container_width=True)
            
            if all_cases:
                st.markdown('<div class="section-title">🏆 经典案例精选</div>', unsafe_allow_html=True)
                
                if len(all_cases) > 3:
                    sample_cases = random.sample(all_cases, 3)
                else:
                    sample_cases = all_cases
                
                for person, case in sample_cases:
                    icon = get_case_icon(case.get('type', '项目案例'))
                    
                    case_html = f"""
                    <div class="case-card">
                        <div class="case-title">
                            {icon} {case.get('title', '经典案例')}
                            <span style="margin-left: auto; color: #6B7280; font-size: 0.9rem;">👤 {person}</span>
                        </div>
                        <div class="case-content">
                            {case.get('content', '')[:100]}...
                        </div>
                    </div>
                    """
                    
                    st.markdown(case_html, unsafe_allow_html=True)
        
        with tab2:
            st.markdown('<div class="section-title">👤 个人档案详情</div>', unsafe_allow_html=True)
            
            col_select, col_actions = st.columns([3, 1])
            with col_select:
                selected_person = st.selectbox("选择客户经理", personnel_list, key="person_select")
            with col_actions:
                if st.button("📋 生成个人报告", key="personal_report"):
                    st.success(f"已生成 {selected_person} 的个人报告")
            
            if selected_person:
                st.markdown("#### 📸 证件照管理")
                handle_photo_upload(selected_person)
                
                st.markdown("""
                <div class="person-card">
                """, unsafe_allow_html=True)
                
                col_info, col_stats, col_photo = st.columns([2, 1, 1])
                
                with col_info:
                    st.markdown(f"### {selected_person}")
                    
                    if selected_person in personnel_info:
                        info = personnel_info[selected_person]
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            if info.get('岗位'):
                                st.markdown(f"**岗位:** {info['岗位']}")
                            if info.get('职称'):
                                st.markdown(f"**职称:** {info['职称']}")
                            if info.get('岗位年限'):
                                st.markdown(f"**岗位年限:** {info['岗位年限']}年")
                        
                        with col2:
                            if info.get('工号'):
                                st.markdown(f"**工号:** {info['工号']}")
                            if info.get('对接客户'):
                                st.markdown(f"**对接客户:** {info['对接客户']}")
                    
                    st.markdown("#### 🏷️ 能力标签")
                    if selected_person in personnel_info and personnel_info[selected_person]['能力标签']:
                        display_skill_tags(personnel_info[selected_person]['能力标签'])
                    else:
                        st.info("暂无能力标签")
                
                with col_stats:
                    completion = position_training_completion[selected_person]
                    
                    st.markdown("#### 📊 培训进度")
                    st.metric("岗位培训完成度", f"{completion['overall_rate']:.1f}%")
                    
                    st.progress(min(completion['overall_rate'] / 100, 1.0))
                    
                    st.write(f"已完成: {completion['total_completed']}/{completion['total_required']} 门必修课")
                
                with col_photo:
                    st.markdown("#### 👤 照片预览")
                    if selected_person in st.session_state.personnel_photos:
                        st.image(st.session_state.personnel_photos[selected_person], 
                                caption=f"{selected_person}",
                                width=150)
                    else:
                        default_img = get_default_photo()
                        st.image(default_img, caption="默认头像", width=150)
                
                st.markdown("</div>", unsafe_allow_html=True)
                
                if selected_person in capability_scores:
                    threshold = st.session_state.warning_thresholds.get(
                        personnel_info[selected_person].get('岗位', '其他'), 70
                    )
                    weaknesses = identify_capability_weaknesses(selected_person, capability_scores, threshold)
                    display_capability_weaknesses(weaknesses)
                
                if selected_person in personnel_info and personnel_info[selected_person]['经典案例']:
                    st.markdown("#### 🏆 岗位经典案例")
                    display_classic_cases(personnel_info[selected_person]['经典案例'])
                
                st.markdown("#### 📈 能力评估雷达图")
                
                if selected_person in capability_scores:
                    radar_fig = create_radar_chart(selected_person, capability_scores)
                    if radar_fig:
                        st.plotly_chart(radar_fig, use_container_width=True)
                        
                        with st.expander("📋 查看详细分数"):
                            col_prof, col_core = st.columns(2)
                            
                            with col_prof:
                                st.write("**五大专业能力得分:**")
                                for category, score in capability_scores[selected_person]['专业能力'].items():
                                    st.write(f"• {category}: {score:.1f}分")
                            
                            with col_core:
                                st.write("**五大核心能力得分:**")
                                for category, score in capability_scores[selected_person]['核心能力'].items():
                                    short_name = category[:10] + "..." if len(category) > 10 else category
                                    st.write(f"• {short_name}: {score:.1f}分")
                    else:
                        st.warning("能力数据不完整")
                else:
                    st.warning("未找到能力评估数据")
        
        with tab3:
            st.markdown('<div class="section-title">🔄 人员对比分析</div>', unsafe_allow_html=True)
            
            col_p1, col_p2 = st.columns(2)
            with col_p1:
                person1 = st.selectbox("选择第一人", personnel_list, key="person1")
            with col_p2:
                other_persons = [p for p in personnel_list if p != person1]
                person2 = st.selectbox("选择第二人", other_persons, key="person2", index=min(1, len(other_persons)-1))
            
            if person1 and person2:
                col_comp1, col_comp2 = st.columns(2)
                
                with col_comp1:
                    st.markdown(f"""
                    <div class="person-card">
                        <h3 style="margin-top: 0;">{person1}</h3>
                        <div style="margin: 15px 0;">
                            <strong>能力标签:</strong>
                    """, unsafe_allow_html=True)
                    
                    if person1 in personnel_info and personnel_info[person1]['能力标签']:
                        display_skill_tags(personnel_info[person1]['能力标签'])
                    else:
                        st.info("暂无能力标签")
                    
                    if person1 in st.session_state.personnel_photos:
                        st.image(st.session_state.personnel_photos[person1], 
                                caption=person1,
                                width=100)
                    
                    st.markdown("""
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                
                with col_comp2:
                    st.markdown(f"""
                    <div class="person-card">
                        <h3 style="margin-top: 0;">{person2}</h3>
                        <div style="margin: 15px 0;">
                            <strong>能力标签:</strong>
                    """, unsafe_allow_html=True)
                    
                    if person2 in personnel_info and personnel_info[person2]['能力标签']:
                        display_skill_tags(personnel_info[person2]['能力标签'])
                    else:
                        st.info("暂无能力标签")
                    
                    if person2 in st.session_state.personnel_photos:
                        st.image(st.session_state.personnel_photos[person2], 
                                caption=person2,
                                width=100)
                    
                    st.markdown("""
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                
                st.markdown("#### 📊 能力对比雷达图")
                
                if person1 in capability_scores and person2 in capability_scores:
                    comparison_fig = create_comparison_radar_chart(person1, person2, capability_scores)
                    if comparison_fig:
                        st.plotly_chart(comparison_fig, use_container_width=True)
                    
                    st.markdown("#### 📋 详细分数对比")
                    
                    if person1 in capability_scores and person2 in capability_scores:
                        col_prof1, col_prof2 = st.columns(2)
                        
                        with col_prof1:
                            st.markdown(f"**{person1} - 专业能力得分:**")
                            prof_scores1 = capability_scores[person1]['专业能力']
                            for category, score in prof_scores1.items():
                                st.write(f"• {category}: {score:.1f}分")
                        
                        with col_prof2:
                            st.markdown(f"**{person2} - 专业能力得分:**")
                            prof_scores2 = capability_scores[person2]['专业能力']
                            for category, score in prof_scores2.items():
                                st.write(f"• {category}: {score:.1f}分")
                        
                        st.markdown("---")
                        col_core1, col_core2 = st.columns(2)
                        
                        with col_core1:
                            st.markdown(f"**{person1} - 核心能力得分:**")
                            core_scores1 = capability_scores[person1]['核心能力']
                            for category, score in core_scores1.items():
                                short_name = category[:10] + "..." if len(category) > 10 else category
                                st.write(f"• {short_name}: {score:.1f}分")
                        
                        with col_core2:
                            st.markdown(f"**{person2} - 核心能力得分:**")
                            core_scores2 = capability_scores[person2]['核心能力']
                            for category, score in core_scores2.items():
                                short_name = category[:10] + "..." if len(category) > 10 else category
                                st.write(f"• {short_name}: {score:.1f}分")
                else:
                    st.warning("无法生成对比雷达图：能力数据不完整")
                
                st.markdown("#### 📚 培训进度对比")
                col_train1, col_train2 = st.columns(2)
                
                with col_train1:
                    completion1 = position_training_completion[person1]
                    st.metric(f"{person1} 培训完成度", f"{completion1['overall_rate']:.1f}%")
                    st.progress(min(completion1['overall_rate'] / 100, 1.0))
                    st.write(f"已完成: {completion1['total_completed']}/{completion1['total_required']} 门必修课")
                
                with col_train2:
                    completion2 = position_training_completion[person2]
                    st.metric(f"{person2} 培训完成度", f"{completion2['overall_rate']:.1f}%")
                    st.progress(min(completion2['overall_rate'] / 100, 1.0))
                    st.write(f"已完成: {completion2['total_completed']}/{completion2['total_required']} 门必修课")
        
        with tab4:
            st.markdown('<div class="section-title">👥 团队综合分析</div>', unsafe_allow_html=True)
            
            st.markdown(f"""
            <div class="threshold-card">
                <h4 style="margin: 0; color: #1E3A8A;">📊 当前团队分析阈值: {st.session_state.team_threshold}分</h4>
                <p style="margin: 5px 0 0 0; color: #666;">此阈值独立于个人预警设置，用于团队能力短板分析</p>
            </div>
            """, unsafe_allow_html=True)
            
            st.markdown('<div class="subsection-title">📊 团队概览</div>', unsafe_allow_html=True)
            
            col1, col2, col3, col4 = st.columns(4)
            
            with col1:
                st.metric("团队总人数", f"{team_report['total_people']}人")
            
            with col2:
                st.metric("平均培训完成度", f"{team_report['avg_completion_rate']:.1f}%")
            
            with col3:
                st.metric("经典案例总数", f"{team_report['total_cases']}个")
            
            with col4:
                st.metric("独特能力标签", f"{team_report['unique_skill_tags_count']}种")
            
            st.markdown('<div class="subsection-title">📚 培训完成情况分析（按岗位要求）</div>', unsafe_allow_html=True)
            
            completion_groups = {}
            for group, persons in team_report['completion_groups_with_names'].items():
                completion_groups[group] = len(persons)
            
            if completion_groups:
                completion_chart = create_training_completion_chart(completion_groups, team_report['completion_groups_with_names'])
                st.plotly_chart(completion_chart, use_container_width=True)
                st.caption("💡 鼠标悬停在柱状图上可查看具体人员名单")
            
            st.markdown('<div class="subsection-title">🎯 团队能力短板分析</div>', unsafe_allow_html=True)
            display_team_weaknesses_with_threshold(team_report['team_weaknesses'], st.session_state.team_threshold)
            
            col_pos, col_title = st.columns(2)
            
            with col_pos:
                st.markdown('<div class="subsection-title">👥 岗位分布分析</div>', unsafe_allow_html=True)
                if team_report['positions']:
                    position_chart = create_position_distribution_chart(team_report['positions'])
                    st.plotly_chart(position_chart, use_container_width=True)
                    
                    position_df = pd.DataFrame({
                        '岗位': list(team_report['positions'].keys()),
                        '人数': list(team_report['positions'].values()),
                        '占比': [f"{(count/team_report['total_people']*100):.1f}%" 
                               for count in team_report['positions'].values()]
                    })
                    st.dataframe(position_df, use_container_width=True)
            
            with col_title:
                st.markdown('<div class="subsection-title">📋 职称分布分析</div>', unsafe_allow_html=True)
                if team_report['titles']:
                    title_chart = create_title_distribution_chart(team_report['titles'])
                    st.plotly_chart(title_chart, use_container_width=True)
                    
                    title_df = pd.DataFrame({
                        '职称': list(team_report['titles'].keys()),
                        '人数': list(team_report['titles'].values()),
                        '占比': [f"{(count/team_report['total_people']*100):.1f}%" 
                               for count in team_report['titles'].values()]
                    })
                    st.dataframe(title_df, use_container_width=True)
            
            st.markdown('<div class="subsection-title">📊 团队能力分析</div>', unsafe_allow_html=True)
            
            if team_report['avg_prof_scores'] or team_report['avg_core_scores']:
                capability_chart = create_capability_comparison_chart(
                    team_report['avg_prof_scores'], team_report['avg_core_scores'],
                    st.session_state.team_threshold
                )
                st.plotly_chart(capability_chart, use_container_width=True)
                st.caption(f"💡 红色虚线为当前团队阈值线: {st.session_state.team_threshold}分")
                
                col_cap1, col_cap2 = st.columns(2)
                
                with col_cap1:
                    if team_report['avg_prof_scores']:
                        st.markdown("##### 📊 专业能力平均得分")
                        prof_df = pd.DataFrame({
                            '能力类别': list(team_report['avg_prof_scores'].keys()),
                            '平均得分': [f"{score:.1f}" for score in team_report['avg_prof_scores'].values()]
                        })
                        st.dataframe(prof_df, use_container_width=True)
                
                with col_cap2:
                    if team_report['avg_core_scores']:
                        st.markdown("##### 💼 核心能力平均得分")
                        core_df = pd.DataFrame({
                            '能力类别': list(team_report['avg_core_scores'].keys()),
                            '平均得分': [f"{score:.1f}" for score in team_report['avg_core_scores'].values()]
                        })
                        st.dataframe(core_df, use_container_width=True)
            
            st.markdown('<div class="subsection-title">🏷️ 能力标签分析</div>', unsafe_allow_html=True)
            if team_report['skill_tag_counts']:
                skill_chart = create_skill_tag_bar_chart(team_report['skill_tag_counts'], top_n=10)
                if skill_chart:
                    st.plotly_chart(skill_chart, use_container_width=True)
                
                sorted_skill_tags = sorted(team_report['skill_tag_counts'].items(), 
                                         key=lambda x: x[1], reverse=True)[:10]
                
                st.markdown("##### 🔥 热门能力标签TOP 10")
                for i, (tag, count) in enumerate(sorted_skill_tags[:10], 1):
                    st.markdown(f"**{i}. {tag}** - {count}人掌握")
            
            st.markdown('<div class="subsection-title">🏆 经典案例统计</div>', unsafe_allow_html=True)
            
            col_case1, col_case2, col_case3 = st.columns(3)
            
            with col_case1:
                st.metric("案例总数", f"{team_report['total_cases']}个")
            
            with col_case2:
                avg_cases_per_person = team_report['total_cases'] / team_report['total_people'] if team_report['total_people'] > 0 else 0
                st.metric("人均案例数", f"{avg_cases_per_person:.1f}个")
            
            with col_case3:
                if all_cases:
                    case_counts = {}
                    for person, _ in all_cases:
                        case_counts[person] = case_counts.get(person, 0) + 1
                    if case_counts:
                        top_case_person = max(case_counts.items(), key=lambda x: x[1])
                        st.metric("案例最多", f"{top_case_person[0]} ({top_case_person[1]}个)")
            
            if team_report['case_years']:
                st.markdown("##### 📅 案例年份分布")
                years_df = pd.DataFrame({
                    '年份': list(team_report['case_years'].keys()),
                    '案例数': list(team_report['case_years'].values())
                })
                years_df = years_df.sort_values('年份', ascending=False)
                st.dataframe(years_df, use_container_width=True)
        
        with tab5:
            st.markdown('<div class="section-title">📋 数据报告与导出</div>', unsafe_allow_html=True)
            
            st.markdown("#### 📊 团队报告生成")
            
            col1, col2 = st.columns([2, 1])
            
            with col1:
                report_name = st.text_input("报告名称", 
                                           value=f"客户经理团队分析报告_{datetime.now().strftime('%Y%m%d')}",
                                           help="请输入报告名称",
                                           key="report_name")
                
                include_warnings = st.checkbox("包含预警信息", value=True, key="include_warnings")
                include_details = st.checkbox("包含详细数据", value=True, key="include_details")
                include_suggestions = st.checkbox("包含改进建议", value=True, key="include_suggestions")
            
            with col2:
                st.markdown("### 🚀")
                if st.button("📄 生成完整报告", type="primary", use_container_width=True, key="generate_report"):
                    with st.spinner("正在生成报告..."):
                        report_content = generate_team_report(
                            personnel_list, personnel_info, training_stats, 
                            capability_scores, all_cases, 
                            warnings_list if include_warnings else [],
                            position_training_completion
                        )
                        
                        st.markdown("### 📄 报告预览")
                        with st.expander("查看报告内容"):
                            st.markdown(report_content)
                        
                        report_bytes = report_content.encode('utf-8')
                        st.download_button(
                            label="💾 下载报告 (TXT)",
                            data=report_bytes,
                            file_name=f"{report_name}.txt",
                            mime="text/plain",
                            key="download_txt"
                        )
                        st.success("✅ 报告生成成功！")
            
            st.markdown("#### 📝 Excel模板使用说明")
            with st.expander("查看模板说明"):
                st.markdown("""
                **Excel文件必须包含以下四个工作表：**
                
                **1. 人员信息表**
                | 姓名 | 工号 | 岗位 | 岗位年限 | 职称 | 对接客户 | 入职日期 | 联系方式 | 能力标签 | 经典案例 |
                |------|------|------|----------|------|----------|----------|----------|----------|----------|
                
                **2. 培训课程表**
                | 培训大类 | 培训课程 | 张三 | 李四 | ... |
                |----------|----------|------|------|-----|
                
                **3. 能力评估表**
                | 姓名 | 通识类 | 商务类 | 业务类 | 产品类 | 管理类 | 客户价值挖掘 | 客户需求捕获 | 客户关系维护 | 投诉抱怨化解 | 资源统筹能力 |
                |------|--------|--------|--------|--------|--------|--------------|--------------|--------------|--------------|--------------|
                
                **4. 岗位培训要求表**
                | 岗位 | 培训课程总数 |
                |------|--------------|
                | 高级客户经理 | 14 |
                | 客户经理 | 10 |
                | 客户经理助理 | 7 |
                
                **说明：**
                - 能力标签列：用逗号分隔多个标签，例如：`英语专八,沟通能力强,数据分析`
                - 经典案例列：使用竖线分隔，格式为`标题|内容|年份|类型|角色|成果`
                - 岗位培训要求表：只需填写岗位名称和培训课程总数两列
                """)
            
            st.markdown("#### 📊 当前数据预览")
            with st.expander("查看原始数据"):
                st.write("**人员信息表:**")
                st.dataframe(data_dict['人员信息'], use_container_width=True)
                
                st.write("**培训课程表:**")
                st.dataframe(data_dict['培训课程'], use_container_width=True)
                
                st.write("**能力评估表:**")
                st.dataframe(data_dict['能力评估'], use_container_width=True)
                
                st.write("**岗位培训要求表:**")
                st.dataframe(data_dict['岗位培训要求'], use_container_width=True)
        
        # ========== 导出功能区域 ==========
        st.markdown("---")
        st.markdown('<div class="section-title">📤 导出功能</div>', unsafe_allow_html=True)
        
        app_url = "http://localhost:8501"
        
        col_status1, col_status2, col_status3 = st.columns(3)
        
        with col_status1:
            if REPORTLAB_AVAILABLE:
                st.success("✅ PDF处理库已安装")
            else:
                st.warning("⚠️ PDF处理库未安装")
        
        with col_status2:
            if PPTX_AVAILABLE:
                st.success("✅ PPT处理库已安装")
            else:
                st.warning("⚠️ PPT处理库未安装")
        
        with col_status3:
            if SELENIUM_AVAILABLE:
                available_browsers = get_available_browsers()
                if available_browsers:
                    st.success("✅ 截图导出库已安装")
                    browser_names = [b[0] for b in available_browsers]
                    st.caption(f"可用浏览器: {', '.join(browser_names)}")
                else:
                    st.warning("⚠️ 未检测到支持的浏览器")
            else:
                st.warning("⚠️ 截图导出库未安装")
        
        st.markdown("---")
        
        export_option = st.radio(
            "选择导出方式",
            ["📸 截图PDF导出（完美保留所有图表和样式）", "📄 标准PDF导出（仅文本）", "📊 PPT导出（仅文本）", "📋 Excel导出（原始数据）"],
            horizontal=True,
            key="export_option"
        )
        
        # 截图导出说明和诊断
        if export_option.startswith("📸") and SELENIUM_AVAILABLE:
            with st.expander("📌 截图导出说明", expanded=False):
                st.info("""
                **使用说明：**
                1. 请确保本应用正在运行（当前窗口就是）
                2. 如果使用Chrome/Edge，需要安装对应的浏览器
                3. 第一次使用时会自动下载浏览器驱动，请耐心等待
                4. 截图过程可能需要10-30秒，请勿关闭页面
                
                **应用访问地址:** http://localhost:8501
                """)
                
                if st.button("🔍 诊断网络连接", key="diagnose_network"):
                    with st.spinner("正在诊断..."):
                        results = diagnose_connection()
                        st.markdown("### 诊断结果：")
                        for r in results:
                            st.write(r)
                        
                        if all("✅" in r for r in results):
                            st.success("✅ 网络连接正常，可以尝试截图导出")
                        else:
                            st.warning("⚠️ 部分地址无法连接，请检查：\n1. 应用是否正常运行\n2. 防火墙是否阻止了连接\n3. 端口8501是否被占用")
        
        selected_browser = 'chrome'
        available_browsers = []
        
        if export_option.startswith("📸") and SELENIUM_AVAILABLE:
            available_browsers = get_available_browsers()
            
            if available_browsers:
                st.markdown('<div class="browser-selector">', unsafe_allow_html=True)
                browser_choice = st.selectbox(
                    "选择浏览器（将自动下载对应的驱动）",
                    options=[b[0] for b in available_browsers],
                    index=0,
                    help="选择用于截图的浏览器，系统会自动下载对应的驱动程序",
                    key="browser_selector"
                )
                browser_map = {b[0]: b[1] for b in available_browsers}
                selected_browser = browser_map[browser_choice]
                st.markdown('</div>', unsafe_allow_html=True)
            else:
                st.error("❌ 未检测到支持的浏览器，请安装 Chrome、Firefox 或 Edge")
                st.info("支持的浏览器：Chrome、Firefox、Microsoft Edge")
        
        if export_option.startswith("📸") and (not SELENIUM_AVAILABLE or not available_browsers):
            st.markdown("""
            <div class="install-box">
                <h4>📥 截图导出需要安装以下组件：</h4>
                <p>1. <b>浏览器</b>（Chrome、Firefox 或 Edge）- 请确保已安装</p>
                <p>2. <b>Python库</b>：在终端运行以下命令：</p>
            </div>
            """, unsafe_allow_html=True)
            
            st.code("pip install selenium webdriver-manager")
            
            if st.button("点击安装所需库", key="install_selenium_btn"):
                install_selenium()
            
            st.stop()
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            if st.button("开始导出", use_container_width=True, type="primary", key="export_button"):
                
                if export_option.startswith("📸"):
                    with st.spinner("正在准备截图导出..."):
                        
                        progress_bar = st.progress(0)
                        status_text = st.empty()
                        
                        def update_progress(message):
                            status_text.text(message)
                            if "完成" in message:
                                progress_bar.progress(100)
                        
                        try:
                            pdf_bytes = create_pdf_from_screenshots(
                                app_url, 
                                browser_type=selected_browser,
                                progress_callback=update_progress
                            )
                            
                            if pdf_bytes:
                                progress_bar.progress(100)
                                status_text.text("✅ 导出完成！")
                                
                                st.download_button(
                                    label="💾 点击下载截图PDF",
                                    data=pdf_bytes,
                                    file_name=f"培训档案截图_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                                    mime="application/pdf",
                                    key="download_screenshot_pdf"
                                )
                                
                        except Exception as e:
                            st.error(f"截图导出失败: {str(e)}")
                            st.info(f"请确保：\n1. {selected_browser.capitalize()}浏览器已安装\n2. 应用正在运行 ({app_url})\n3. 网络连接正常")
                            st.info("如果是第一次使用，可能需要等待驱动程序下载完成")
                
                elif export_option == "📄 标准PDF导出（仅文本）":
                    if not REPORTLAB_AVAILABLE:
                        st.error("PDF导出功能不可用：缺少reportlab库")
                        st.code("pip install reportlab")
                    else:
                        with st.spinner("正在生成PDF文件..."):
                            pdf_buffer = export_all_to_pdf(
                                personnel_list, personnel_info, capability_scores, training_stats,
                                position_training_completion, all_cases, warnings_list, team_report
                            )
                            if pdf_buffer:
                                st.download_button(
                                    label="💾 点击下载PDF文件",
                                    data=pdf_buffer,
                                    file_name=f"培训档案报告_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                                    mime="application/pdf",
                                    key="download_pdf"
                                )
                
                elif export_option == "📊 PPT导出（仅文本）":
                    if not PPTX_AVAILABLE:
                        st.error("PPT导出功能不可用：缺少python-pptx库")
                        st.code("pip install python-pptx")
                    else:
                        with st.spinner("正在生成PPT文件..."):
                            ppt_buffer = export_all_to_ppt(
                                personnel_list, personnel_info, capability_scores, training_stats,
                                position_training_completion, all_cases, warnings_list, team_report
                            )
                            if ppt_buffer:
                                st.download_button(
                                    label="💾 点击下载PPT文件",
                                    data=ppt_buffer,
                                    file_name=f"培训档案报告_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key="download_ppt"
                                )
                
                elif export_option == "📋 Excel导出（原始数据）":
                    with st.spinner("正在生成Excel文件..."):
                        try:
                            output = io.BytesIO()
                            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                                if '人员信息' in data_dict:
                                    data_dict['人员信息'].to_excel(writer, sheet_name='人员信息', index=False)
                                if '培训课程' in data_dict:
                                    data_dict['培训课程'].to_excel(writer, sheet_name='培训课程', index=False)
                                if '能力评估' in data_dict:
                                    data_dict['能力评估'].to_excel(writer, sheet_name='能力评估', index=False)
                                if '岗位培训要求' in data_dict:
                                    data_dict['岗位培训要求'].to_excel(writer, sheet_name='岗位培训要求', index=False)
                                
                                if team_report:
                                    if team_report['positions']:
                                        pos_df = pd.DataFrame({
                                            '岗位': list(team_report['positions'].keys()),
                                            '人数': list(team_report['positions'].values())
                                        })
                                        pos_df.to_excel(writer, sheet_name='岗位分布', index=False)
                                    
                                    if team_report['skill_tag_counts']:
                                        tag_df = pd.DataFrame({
                                            '能力标签': list(team_report['skill_tag_counts'].keys()),
                                            '掌握人数': list(team_report['skill_tag_counts'].values())
                                        })
                                        tag_df = tag_df.sort_values('掌握人数', ascending=False)
                                        tag_df.to_excel(writer, sheet_name='热门能力标签', index=False)
                            
                            output.seek(0)
                            st.download_button(
                                label="💾 点击下载Excel文件",
                                data=output,
                                file_name=f"培训档案数据_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                key="download_excel"
                            )
                        except Exception as e:
                            st.error(f"Excel生成失败: {str(e)}")
        
        st.markdown("---")
        st.markdown(f"""
        <div style="text-align: center; color: #666; padding: 20px;">
            <div>📅 数据更新时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</div>
            <div>📁 数据文件: {st.session_state.uploaded_file_name} | 👥 总人数: {len(personnel_list)} | 🏆 总案例数: {len(all_cases)}</div>
            <div style="margin-top: 10px; font-size: 0.9rem;">
                © 2024 客户经理培训档案管理系统
            </div>
        </div>
        """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()